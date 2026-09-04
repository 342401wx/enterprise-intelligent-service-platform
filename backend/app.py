from __future__ import annotations

import hashlib
import html as html_lib
import hmac
import secrets as crypto_secrets
import json
import os
import re
import shutil
import subprocess
import threading
import sqlite3
import time
import uuid
import urllib.error
import urllib.parse
import urllib.request
from contextlib import closing
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Any, Literal

from fastapi import Depends, FastAPI, File, Form, Header, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel, Field
from backend.model_discovery import fetch_models as fetch_models_from_provider
from backend.secrets import decrypt_secret, encrypt_secret


ROOT = Path(__file__).resolve().parent
DATA_DIR = Path(os.getenv("PLATFORM_DATA_DIR", str(ROOT / "data")))
DB_PATH = Path(os.getenv("PLATFORM_DB_PATH", str(DATA_DIR / "platform.db")))
DB_PATH.parent.mkdir(parents=True, exist_ok=True)
EMAIL_ATTACHMENT_DIR = DATA_DIR / 'email-attachments'
EMAIL_ATTACHMENT_DIR.mkdir(parents=True, exist_ok=True)

EMAIL_AUTH_JOBS: dict[str, dict[str, Any]] = {}
EMAIL_AUTH_LOCK = threading.Lock()


def now() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


def connect() -> sqlite3.Connection:
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA foreign_keys = ON")
    return conn


def rows(conn: sqlite3.Connection, query: str, params: tuple[Any, ...] = ()) -> list[dict[str, Any]]:
    return [dict(row) for row in conn.execute(query, params).fetchall()]


def row(conn: sqlite3.Connection, query: str, params: tuple[Any, ...] = ()) -> dict[str, Any] | None:
    item = conn.execute(query, params).fetchone()
    return dict(item) if item else None


def audit(conn: sqlite3.Connection, user_id: str, action: str, resource_type: str, resource_id: str | None, detail: dict[str, Any] | None = None) -> None:
    conn.execute(
        "INSERT INTO audit_events (id, user_id, action, resource_type, resource_id, detail, created_at) VALUES (?, ?, ?, ?, ?, ?, ?)",
        (f"AUD-{uuid.uuid4().hex[:12].upper()}", user_id, action, resource_type, resource_id, json.dumps(detail or {}, ensure_ascii=False), now()),
    )


def _agentmail_command(arguments: list[str]) -> list[str]:
    candidates = [
        os.getenv('AGENTLY_CLI_PATH', ''),
        shutil.which('agently-cli') or '',
        shutil.which('agently-cli.cmd') or '',
        shutil.which('agently-cli.ps1') or '',
        str(Path(os.getenv('APPDATA', '')) / 'npm' / 'agently-cli.ps1'),
    ]
    cli = next((item for item in candidates if item and Path(item).exists()), None)
    if not cli:
        raise RuntimeError('未找到 AgentMail CLI，请先安装 @tencent-qqmail/agently-cli')
    if cli.lower().endswith('.ps1'):
        return ['powershell.exe', '-NoProfile', '-ExecutionPolicy', 'Bypass', '-File', cli, *arguments]
    return [cli, *arguments]


def _run_agentmail(arguments: list[str], workspace: str | None = None) -> dict[str, Any]:
    env = os.environ.copy()
    if workspace:
        env["AGENTLY_WORKSPACE"] = workspace
    try:
        completed = subprocess.run(
            _agentmail_command(arguments), cwd=ROOT, env=env, capture_output=True, text=True,
            encoding='utf-8', errors='replace', timeout=60, check=False,
        )
    except subprocess.TimeoutExpired as exc:
        raise RuntimeError('AgentMail 请求超时') from exc
    except OSError as exc:
        raise RuntimeError(f'AgentMail CLI 启动失败：{exc}') from exc
    raw = (completed.stdout or '').strip()
    payload: dict[str, Any] | None = None
    decoder = json.JSONDecoder()
    for index, character in enumerate(raw):
        if character != '{':
            continue
        try:
            candidate, _ = decoder.raw_decode(raw[index:])
        except json.JSONDecodeError:
            continue
        if isinstance(candidate, dict):
            payload = candidate
            break
    if completed.returncode != 0 or not payload or not payload.get('ok'):
        error = (payload or {}).get('error') or {}
        message = error.get('message') if isinstance(error, dict) else None
        raise RuntimeError(str(message or raw[-800:] or 'AgentMail 请求失败'))
    return payload.get('data') if isinstance(payload.get('data'), dict) else {}

def _active_account_id(user: dict[str, Any], account_id: str | None = None) -> str | None:
    with closing(connect()) as conn:
        if account_id:
            account = row(conn, "SELECT id FROM email_accounts WHERE id = ? AND owner_id = ? AND status = 'active'", (account_id, user["id"]))
            if not account:
                raise HTTPException(404, "邮箱账号不存在或无权使用")
            return str(account["id"])
        account = row(conn, "SELECT id FROM email_accounts WHERE owner_id = ? AND is_active = 1 AND status = 'active' ORDER BY updated_at DESC LIMIT 1", (user["id"],))
    return str(account["id"]) if account else None
def _account_workspace(user: dict[str, Any], account_id: str | None = None) -> str:
    with closing(connect()) as conn:
        if account_id:
            account = row(conn, "SELECT workspace FROM email_accounts WHERE id = ? AND owner_id = ? AND status = 'active'", (account_id, user["id"]))
            if not account:
                raise HTTPException(404, "邮箱账号不存在或无权使用")
            return None if account["workspace"] == "default" else str(account["workspace"])
        account = row(conn, "SELECT workspace FROM email_accounts WHERE owner_id = ? AND is_active = 1 AND status = 'active' ORDER BY updated_at DESC LIMIT 1", (user["id"],))
    if not account or account["workspace"] == "default": return None
    return str(account["workspace"])


def _account_public(item: dict[str, Any]) -> dict[str, Any]:
    return {"id": item["id"], "workspace": item["workspace"], "email": item["email"], "name": item.get("name") or item["email"], "status": item["status"], "active": bool(item.get("is_active"))}


def _save_agent_account(user_id: str, workspace: str, profile: dict[str, Any], label: str = "", make_active: bool = False) -> dict[str, Any]:
    aliases = profile.get("aliases") if isinstance(profile.get("aliases"), list) else []
    primary = next((item for item in aliases if isinstance(item, dict) and item.get("is_primary")), None)
    primary = primary or next((item for item in aliases if isinstance(item, dict)), {})
    email = str(primary.get("email") or "")
    name = str(primary.get("name") or label or email)
    if not email:
        raise RuntimeError("AgentMail 授权成功，但未返回邮箱地址")
    stamp = now()
    with closing(connect()) as conn:
        existing = row(conn, "SELECT id FROM email_accounts WHERE owner_id = ? AND workspace = ?", (user_id, workspace))
        account_id = existing["id"] if existing else f"EA-{uuid.uuid4().hex[:10].upper()}"
        active_exists = row(conn, "SELECT id FROM email_accounts WHERE owner_id = ? AND is_active = 1", (user_id,))
        is_active = 1 if make_active or not active_exists or active_exists["id"] == account_id else 0
        conn.execute("INSERT INTO email_accounts (id, owner_id, workspace, email, name, status, is_active, created_at, updated_at) VALUES (?, ?, ?, ?, ?, 'active', ?, ?, ?) ON CONFLICT(owner_id, workspace) DO UPDATE SET email=excluded.email, name=excluded.name, status='active', updated_at=excluded.updated_at", (account_id, user_id, workspace, email, name, is_active, stamp, stamp))
        if make_active:
            conn.execute("UPDATE email_accounts SET is_active = 0 WHERE owner_id = ? AND id <> ?", (user_id, account_id))
        conn.execute("UPDATE email_messages SET account_id = ? WHERE owner_id = ? AND account_id IS NULL", (account_id, user_id))
        conn.commit()
        return _account_public(row(conn, "SELECT * FROM email_accounts WHERE id = ?", (account_id,)) or {})


def _set_auth_job(job_id: str, **changes: Any) -> None:
    with EMAIL_AUTH_LOCK:
        if job_id in EMAIL_AUTH_JOBS:
            EMAIL_AUTH_JOBS[job_id].update(changes)


def _run_email_auth(job_id: str, user_id: str, workspace: str, label: str) -> None:
    env = os.environ.copy()
    env["AGENTLY_WORKSPACE"] = workspace
    process = None
    try:
        process = subprocess.Popen(_agentmail_command(["auth", "login", "--verbose"]), cwd=ROOT, env=env, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True, encoding="utf-8", errors="replace", bufsize=1)
        _set_auth_job(job_id, pid=process.pid)
        output: list[str] = []
        if process.stdout:
            for line in process.stdout:
                clean = line.rstrip()
                output.append(clean)
                match = re.search(r"https?://[^\s]+", clean)
                _set_auth_job(job_id, output="\n".join(output[-40:]), authorization_url=match.group(0) if match else EMAIL_AUTH_JOBS.get(job_id, {}).get("authorization_url"))
        return_code = process.wait(timeout=300)
        if return_code != 0:
            raise RuntimeError("微信扫码授权未完成，请重新发起授权")
        profile = _run_agentmail(["+me"], workspace=workspace)
        account = _save_agent_account(user_id, workspace, profile, label, make_active=True)
        _set_auth_job(job_id, status="completed", account=account, message=f"邮箱 {account['email']} 授权成功")
    except subprocess.TimeoutExpired:
        if process:
            process.kill()
        _set_auth_job(job_id, status="failed", message="授权等待超时，请重新发起微信扫码授权")
    except Exception as exc:
        _set_auth_job(job_id, status="failed", message=str(exc))

def _json_text(value: Any) -> str:
    return json.dumps(value if value is not None else [], ensure_ascii=False)


def _mail_provider_key(provider_id: str, account_id: str | None) -> str:
    return f"{account_id}:{provider_id}" if account_id else provider_id
def _email_address(value: Any) -> dict[str, str]:
    if isinstance(value, dict):
        return {'email': str(value.get('email', '')), 'name': str(value.get('name', ''))}
    return {'email': str(value or ''), 'name': ''}


def _upsert_email(conn: sqlite3.Connection, owner_id: str, message: dict[str, Any], folder: str, source: str, body: str | None = None, account_id: str | None = None) -> dict[str, Any]:
    remote_provider_id = str(message.get('message_id') or message.get('id') or f'msg_local_{uuid.uuid4().hex}')
    provider_id = _mail_provider_key(remote_provider_id, account_id)
    sender = _email_address(message.get('from'))
    recipients = [_email_address(item) for item in (message.get('to') or [])]
    cc = [_email_address(item) for item in (message.get('cc') or [])]
    attachments = message.get('attachments') or []
    created_at = str(message.get('created_at') or now())
    content = str(body or message.get('body') or '')
    # Reading a full message may return HTML without a snippet. Keep the existing list preview in that case instead of replacing it with markup.
    preview = str(message.get('snippet') or (content[:180] if source != 'read' else ''))
    direction = 'sent' if folder == 'sent' else 'received'
    status = 'sent' if direction == 'sent' else 'received'
    conn.execute(
        '''INSERT INTO email_messages
        (id, provider_id, owner_id, account_id, folder, direction, sender_email, sender_name, recipients, cc,
         subject, preview, body, is_read, has_attachments, attachments, source, status,
         provider_payload, created_at, updated_at)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        ON CONFLICT(provider_id) DO UPDATE SET
          folder=excluded.folder, direction=excluded.direction, sender_email=excluded.sender_email,
          sender_name=excluded.sender_name, recipients=excluded.recipients, cc=excluded.cc,
          subject=excluded.subject,
          preview=CASE WHEN excluded.preview <> '' THEN excluded.preview ELSE email_messages.preview END,
          body=CASE WHEN excluded.body <> '' THEN excluded.body ELSE email_messages.body END,
          is_read=CASE WHEN excluded.is_read = 1 THEN 1 ELSE email_messages.is_read END,
          has_attachments=excluded.has_attachments, attachments=excluded.attachments,
          source=excluded.source, status=excluded.status, account_id=excluded.account_id, provider_payload=excluded.provider_payload,
          updated_at=excluded.updated_at''',
        (f'MAIL-{uuid.uuid4().hex[:12].upper()}', provider_id, owner_id, account_id, folder, direction,
         sender['email'], sender['name'], _json_text(recipients), _json_text(cc),
         str(message.get('subject') or ''), preview, content, int(bool(message.get('is_read', direction == 'sent'))),
         int(bool(message.get('has_attachments') or attachments)), _json_text(attachments), source, status,
         json.dumps(message, ensure_ascii=False), created_at, now()),
    )
    return row(conn, 'SELECT * FROM email_messages WHERE provider_id = ?', (provider_id,)) or {}


def _email_response(item: dict[str, Any]) -> dict[str, Any]:
    result = dict(item)
    for field in ('recipients', 'cc', 'attachments'):
        try:
            result[field] = json.loads(result.get(field) or '[]')
        except json.JSONDecodeError:
            result[field] = []
    result['unread'] = not bool(result.pop('is_read', 0))
    stored_provider_id = result.pop('provider_id', '')
    try:
        provider_payload = json.loads(result.get('provider_payload') or '{}')
    except json.JSONDecodeError:
        provider_payload = {}
    result.pop('provider_payload', None)
    result['providerId'] = str(provider_payload.get('message_id') or provider_payload.get('id') or stored_provider_id.split(':', 1)[-1])
    result['from'] = result.pop('sender_email', '')
    result['fromName'] = result.pop('sender_name', '')
    result['to'] = result.pop('recipients', [])
    result['createdAt'] = result.pop('created_at', '')
    result['updatedAt'] = result.pop('updated_at', '')
    result.pop('owner_id', None)
    result.pop('provider_payload', None)
    return result


PASSWORD_ITERATIONS = 260_000
TEST_PASSWORD = os.getenv("PLATFORM_TEST_PASSWORD", "12345678")


def hash_password(password: str) -> str:
    if len(password) < 8:
        raise ValueError("密码至少需要 8 位")
    salt = crypto_secrets.token_bytes(16)
    digest = hashlib.pbkdf2_hmac("sha256", password.encode("utf-8"), salt, PASSWORD_ITERATIONS)
    return "pbkdf2_sha256$" + str(PASSWORD_ITERATIONS) + "$" + salt.hex() + "$" + digest.hex()


def hash_test_password(password: str) -> str:
    if len(password) < 6:
        raise ValueError("测试密码至少需要 6 位")
    salt = crypto_secrets.token_bytes(16)
    digest = hashlib.pbkdf2_hmac("sha256", password.encode("utf-8"), salt, PASSWORD_ITERATIONS)
    return "pbkdf2_sha256$" + str(PASSWORD_ITERATIONS) + "$" + salt.hex() + "$" + digest.hex()

def verify_password(password: str, encoded: str) -> bool:
    try:
        algorithm, iterations, salt_hex, digest_hex = encoded.split("$")
        if algorithm != "pbkdf2_sha256":
            return False
        actual = hashlib.pbkdf2_hmac("sha256", password.encode("utf-8"), bytes.fromhex(salt_hex), int(iterations))
        return hmac.compare_digest(actual.hex(), digest_hex)
    except (ValueError, TypeError):
        return False


def public_user(user: dict[str, Any]) -> dict[str, Any]:
    result = dict(user)
    result.pop("password_hash", None)
    return result

def init_db() -> None:
    with closing(connect()) as conn:
        conn.executescript(
            """
            CREATE TABLE IF NOT EXISTS users (
                id TEXT PRIMARY KEY, name TEXT NOT NULL, email TEXT NOT NULL UNIQUE,
                department TEXT NOT NULL, role TEXT NOT NULL, status TEXT NOT NULL DEFAULT 'active', password_hash TEXT NOT NULL DEFAULT '', supervisor_id TEXT
            );
            CREATE TABLE IF NOT EXISTS auth_sessions (
                token_hash TEXT PRIMARY KEY, user_id TEXT NOT NULL, created_at TEXT NOT NULL, expires_at TEXT NOT NULL,
                FOREIGN KEY(user_id) REFERENCES users(id)
            );
            CREATE TABLE IF NOT EXISTS departments (
                id TEXT PRIMARY KEY, name TEXT NOT NULL UNIQUE, parent_id TEXT,
                manager_id TEXT, status TEXT NOT NULL DEFAULT 'active',
                FOREIGN KEY(manager_id) REFERENCES users(id)
            );            CREATE TABLE IF NOT EXISTS conversations (
                id TEXT PRIMARY KEY, user_id TEXT NOT NULL, title TEXT NOT NULL,
                preview TEXT NOT NULL DEFAULT '', status TEXT NOT NULL DEFAULT 'active',
                created_at TEXT NOT NULL, updated_at TEXT NOT NULL,
                FOREIGN KEY(user_id) REFERENCES users(id)
            );
            CREATE TABLE IF NOT EXISTS messages (
                id TEXT PRIMARY KEY, conversation_id TEXT NOT NULL, role TEXT NOT NULL,
                content TEXT NOT NULL, seq INTEGER NOT NULL, created_at TEXT NOT NULL,
                FOREIGN KEY(conversation_id) REFERENCES conversations(id)
            );
            CREATE TABLE IF NOT EXISTS conversation_attachments (
                id TEXT PRIMARY KEY, conversation_id TEXT NOT NULL, message_seq INTEGER NOT NULL,
                attachment_id TEXT NOT NULL, filename TEXT NOT NULL, size INTEGER NOT NULL DEFAULT 0,
                content_type TEXT NOT NULL DEFAULT 'application/octet-stream', provider_message_id TEXT NOT NULL,
                account_id TEXT, created_at TEXT NOT NULL,
                UNIQUE(conversation_id, message_seq, attachment_id),
                FOREIGN KEY(conversation_id) REFERENCES conversations(id)
            );            CREATE TABLE IF NOT EXISTS todo_items (
                id TEXT PRIMARY KEY, owner_id TEXT NOT NULL, title TEXT NOT NULL,
                due_date TEXT, priority TEXT NOT NULL DEFAULT 'normal', status TEXT NOT NULL DEFAULT 'open',
                created_at TEXT NOT NULL, updated_at TEXT NOT NULL,
                FOREIGN KEY(owner_id) REFERENCES users(id)
            );            CREATE TABLE IF NOT EXISTS password_reset_requests (
                id TEXT PRIMARY KEY, user_id TEXT NOT NULL, identifier TEXT NOT NULL,
                status TEXT NOT NULL DEFAULT 'pending', created_at TEXT NOT NULL,
                read_at TEXT, handled_by TEXT, handled_at TEXT,
                FOREIGN KEY(user_id) REFERENCES users(id)
            );            CREATE TABLE IF NOT EXISTS agent_events (
                id TEXT PRIMARY KEY, conversation_id TEXT NOT NULL, task_id TEXT NOT NULL,
                trace_id TEXT NOT NULL, seq INTEGER NOT NULL, event_type TEXT NOT NULL,
                label TEXT NOT NULL, status TEXT NOT NULL, duration TEXT NOT NULL,
                safe_summary TEXT NOT NULL, data TEXT NOT NULL, created_at TEXT NOT NULL,
                FOREIGN KEY(conversation_id) REFERENCES conversations(id)
            );
            CREATE TABLE IF NOT EXISTS token_usage (
                id TEXT PRIMARY KEY, user_id TEXT NOT NULL, conversation_id TEXT NOT NULL,
                task_id TEXT NOT NULL, provider TEXT NOT NULL, model TEXT NOT NULL,
                input_tokens INTEGER NOT NULL DEFAULT 0, output_tokens INTEGER NOT NULL DEFAULT 0,
                total_tokens INTEGER NOT NULL DEFAULT 0, created_at TEXT NOT NULL,
                FOREIGN KEY(user_id) REFERENCES users(id)
            );            CREATE TABLE IF NOT EXISTS leave_requests (
                id TEXT PRIMARY KEY, applicant_id TEXT NOT NULL, applicant TEXT NOT NULL,
                department TEXT NOT NULL, leave_type TEXT NOT NULL, dates TEXT NOT NULL,
                days REAL NOT NULL, status TEXT NOT NULL, reason TEXT NOT NULL,
                approver_id TEXT, created_at TEXT NOT NULL, updated_at TEXT NOT NULL,
                FOREIGN KEY(applicant_id) REFERENCES users(id)
            );
            CREATE TABLE IF NOT EXISTS documents (
                id TEXT PRIMARY KEY, name TEXT NOT NULL, owner TEXT NOT NULL, size TEXT NOT NULL,
                status TEXT NOT NULL, stage TEXT NOT NULL, updated TEXT NOT NULL,
                knowledge_base TEXT NOT NULL, created_at TEXT NOT NULL
            );
            CREATE TABLE IF NOT EXISTS ingestion_tasks (
                id TEXT PRIMARY KEY, document_id TEXT NOT NULL, document TEXT NOT NULL,
                stage TEXT NOT NULL, status TEXT NOT NULL, progress INTEGER NOT NULL,
                updated TEXT NOT NULL
            );
            CREATE TABLE IF NOT EXISTS generated_files (
                id TEXT PRIMARY KEY, name TEXT NOT NULL, file_type TEXT NOT NULL,
                status TEXT NOT NULL, conversation_id TEXT NOT NULL, created_at TEXT NOT NULL,
                template TEXT NOT NULL
            );
            CREATE TABLE IF NOT EXISTS notifications (
                id TEXT PRIMARY KEY, title TEXT NOT NULL, detail TEXT NOT NULL,
                type TEXT NOT NULL, time TEXT NOT NULL, unread INTEGER NOT NULL DEFAULT 1
            );
            CREATE TABLE IF NOT EXISTS model_configs (
                id TEXT PRIMARY KEY, scope_type TEXT NOT NULL, scope_id TEXT NOT NULL,
                provider TEXT NOT NULL, model TEXT NOT NULL, api_url TEXT NOT NULL,
                api_key_masked TEXT NOT NULL, api_format TEXT NOT NULL DEFAULT 'openai-responses',
                api_key_ciphertext TEXT NOT NULL DEFAULT '', enabled INTEGER NOT NULL DEFAULT 1,
                updated_at TEXT NOT NULL, UNIQUE(scope_type, scope_id)
            );
            CREATE TABLE IF NOT EXISTS email_messages (
                id TEXT PRIMARY KEY, provider_id TEXT NOT NULL UNIQUE, owner_id TEXT NOT NULL,
                folder TEXT NOT NULL, direction TEXT NOT NULL, sender_email TEXT NOT NULL,
                sender_name TEXT NOT NULL DEFAULT '', recipients TEXT NOT NULL DEFAULT '[]',
                cc TEXT NOT NULL DEFAULT '[]', subject TEXT NOT NULL DEFAULT '',
                preview TEXT NOT NULL DEFAULT '', body TEXT NOT NULL DEFAULT '',
                is_read INTEGER NOT NULL DEFAULT 0, has_attachments INTEGER NOT NULL DEFAULT 0,
                attachments TEXT NOT NULL DEFAULT '[]', source TEXT NOT NULL DEFAULT 'sync',
                status TEXT NOT NULL DEFAULT 'received', provider_payload TEXT NOT NULL DEFAULT '{}',
                created_at TEXT NOT NULL, updated_at TEXT NOT NULL
            );            CREATE TABLE IF NOT EXISTS email_accounts (
                id TEXT PRIMARY KEY, owner_id TEXT NOT NULL, workspace TEXT NOT NULL,
                email TEXT NOT NULL, name TEXT NOT NULL DEFAULT '', status TEXT NOT NULL DEFAULT 'active',
                is_active INTEGER NOT NULL DEFAULT 0, created_at TEXT NOT NULL, updated_at TEXT NOT NULL,
                UNIQUE(owner_id, workspace)
            );            CREATE TABLE IF NOT EXISTS audit_events (
                id TEXT PRIMARY KEY, user_id TEXT NOT NULL, action TEXT NOT NULL,
                resource_type TEXT NOT NULL, resource_id TEXT, detail TEXT NOT NULL,
                created_at TEXT NOT NULL
            );
            """
        )
        user_columns = {item[1] for item in conn.execute("PRAGMA table_info(users)").fetchall()}
        if "supervisor_id" not in user_columns:
            conn.execute("ALTER TABLE users ADD COLUMN supervisor_id TEXT")
        if "password_hash" not in user_columns:
            conn.execute("ALTER TABLE users ADD COLUMN password_hash TEXT NOT NULL DEFAULT ''")
        email_columns = {item[1] for item in conn.execute("PRAGMA table_info(email_messages)").fetchall()}
        if "account_id" not in email_columns:
            conn.execute("ALTER TABLE email_messages ADD COLUMN account_id TEXT")
        legacy_rows = conn.execute("SELECT rowid, provider_id, account_id, provider_payload FROM email_messages WHERE account_id IS NOT NULL").fetchall()
        for legacy in legacy_rows:
            stored_id = str(legacy["provider_id"] or "")
            account_ref = str(legacy["account_id"] or "")
            if not account_ref or stored_id.startswith(account_ref + ":"):
                continue
            desired_id = f"{account_ref}:{stored_id}"
            conflict = conn.execute("SELECT rowid FROM email_messages WHERE provider_id = ? AND rowid <> ?", (desired_id, legacy["rowid"])).fetchone()
            if conflict:
                desired_id = "{}:legacy:{}:{}".format(account_ref, legacy["rowid"], stored_id)
            conn.execute("UPDATE email_messages SET provider_id = ? WHERE rowid = ?", (desired_id, legacy["rowid"]))

        file_columns = {item[1] for item in conn.execute("PRAGMA table_info(generated_files)").fetchall()}
        if "owner_id" not in file_columns:
            conn.execute("ALTER TABLE generated_files ADD COLUMN owner_id TEXT NOT NULL DEFAULT ''")

        if "file_path" not in file_columns:
            conn.execute("ALTER TABLE generated_files ADD COLUMN file_path TEXT NOT NULL DEFAULT ''")
        columns = {item[1] for item in conn.execute("PRAGMA table_info(model_configs)").fetchall()}
        if "api_format" not in columns:
            conn.execute("ALTER TABLE model_configs ADD COLUMN api_format TEXT NOT NULL DEFAULT 'openai-responses'")
        if "api_key_ciphertext" not in columns:
            conn.execute("ALTER TABLE model_configs ADD COLUMN api_key_ciphertext TEXT NOT NULL DEFAULT ''")
        seed(conn)
        legacy_email_messages = conn.execute("SELECT id, content FROM messages WHERE role = 'assistant' AND (content LIKE '%---原始邮件---%' OR content LIKE '%qmbox%' OR content LIKE '%&lt;style%')").fetchall()
        for legacy in legacy_email_messages:
            content = html_lib.unescape(str(legacy['content'] or ''))
            content = re.sub(r'\.qmbox\b.*?(?=Agent Mail 接入成功|邮箱地址)', ' ', content, flags=re.S | re.I)
            content = re.sub(r'---\s*原始邮件\s*---.*$', '', content, flags=re.S | re.I)
            content = re.sub(r'\s+', ' ', content).strip()
            content = re.sub(r'(第一封[^：]*：)\s*主题：', r'\1\n主题：', content)
            content = content.replace(' 发件人：', '\n发件人：').replace(' 正文：', '\n正文：')
            conn.execute('UPDATE messages SET content = ? WHERE id = ?', (content, legacy['id']))
        conn.commit()


def _fixture_json(name: str) -> Any:
    fixture = ROOT.parent / 'test-fixtures' / 'organization' / name
    if not fixture.exists():
        raise RuntimeError(f'测试数据文件不存在：{fixture}')
    return json.loads(fixture.read_text(encoding='utf-8'))


def seed(conn: sqlite3.Connection) -> None:
    users = _fixture_json('users.json')
    for item in users:
        existing = conn.execute("SELECT id FROM users WHERE email = ?", (item["email"],)).fetchone()
        password_hash = hash_test_password(item.get("initial_password", TEST_PASSWORD))
        if existing:
            conn.execute(
                "UPDATE users SET name = ?, department = ?, role = ?, status = ?, password_hash = ?, supervisor_id = ? WHERE email = ?",
                (item["name"], item["department"], item["role"], item.get("status", "active"), password_hash, item.get("supervisor_id"), item["email"]),
            )
        else:
            conn.execute(
                "INSERT INTO users (id, name, email, department, role, status, password_hash, supervisor_id) VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
                (item["id"], item["name"], item["email"], item["department"], item["role"], item.get("status", "active"), password_hash, item.get("supervisor_id")),
            )
    departments = _fixture_json('departments.json')
    conn.executemany(
        'INSERT OR IGNORE INTO departments (id, name, parent_id, manager_id, status) VALUES (?, ?, ?, ?, ?)',
        [(item['id'], item['name'], item.get('parent_id'), item.get('manager_id'), item.get('status', 'active')) for item in departments],
    )
    stamp = now()

    knowledge_root = ROOT.parent / 'test-fixtures' / 'knowledge'
    knowledge_base = {
        '员工手册与年假制度.md': '企业制度库',
        '财务管理制度.md': '财务制度库',
        '考勤与办公管理制度.md': '企业制度库',
        '信息安全与文档管理制度.md': '信息安全库',
        '采购与合同管理制度.md': '采购制度库',
    }
    for name, knowledge_base_name in knowledge_base.items():
        source = knowledge_root / name
        if not source.exists():
            raise RuntimeError(f'知识库测试文件不存在：{source}')
        size = f'{source.stat().st_size / 1024:.1f} KB'
        document_id = 'DOC-' + uuid.uuid5(uuid.NAMESPACE_URL, 'enterprise-fixture:' + name).hex[:10].upper()
        task_id = 'ING-' + uuid.uuid5(uuid.NAMESPACE_URL, 'enterprise-ingestion:' + name).hex[:10].upper()
        conn.execute(
            'INSERT OR IGNORE INTO documents VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)',
            (document_id, name, '测试知识库管理员', size, 'success', '已发布', '测试数据', knowledge_base_name, stamp),
        )
        conn.execute(
            'INSERT OR IGNORE INTO ingestion_tasks VALUES (?, ?, ?, ?, ?, ?, ?)',
            (task_id, document_id, name, '已完成', 'success', 100, '测试数据'),
        )

class LoginRequest(BaseModel):
    email: str
    password: str


class ForgotPasswordRequest(BaseModel):
    identifier: str = Field(min_length=1, max_length=160)

class AdminPasswordResetRequest(BaseModel):
    new_password: str = Field(min_length=8, max_length=128)
    request_id: str | None = None
class CreateUserRequest(BaseModel):
    name: str = Field(min_length=1, max_length=80)
    email: str
    department: str = Field(min_length=1, max_length=80)
    role: Literal["employee", "manager", "admin"] = "employee"
    password: str = Field(min_length=8, max_length=128)

class LeaveCreate(BaseModel):
    leave_type: str = Field(alias="type")
    start: str
    end: str
    reason: str
    days: float = 3

    model_config = {"populate_by_name": True}


class DecisionRequest(BaseModel):
    action: Literal["approved", "rejected", "returned"]
    comment: str = ""


class TodoCreateRequest(BaseModel):
    title: str = Field(min_length=1, max_length=160)
    due_date: str | None = None
    priority: Literal["high", "normal", "low"] = "normal"

class TodoUpdateRequest(BaseModel):
    status: Literal["open", "done"]
class MessageRequest(BaseModel):
    content: str


class ModelConfigRequest(BaseModel):
    scope_type: Literal["global", "independent"] = "global"
    scope_id: str = "platform"
    provider: str
    model: str
    api_url: str
    api_key: str = ""
    api_format: str = "openai-responses"
    enabled: bool = True


class EmailAuthStartRequest(BaseModel):
    label: str = Field(min_length=1, max_length=80)
    workspace: str | None = None

class AgentMailSendRequest(BaseModel):
    to: list[str] = Field(min_length=1)
    subject: str = ''
    body: str = ''
    cc: list[str] = Field(default_factory=list)
    bcc: list[str] = Field(default_factory=list)
    confirmed: bool = False
    account_id: str | None = None

class ModelDiscoveryRequest(BaseModel):
    provider: str = "OpenAI Compatible"
    api_url: str
    api_key: str
    api_format: str = "openai-responses"
    is_full_url: bool = False
    models_url_override: str | None = None


app = FastAPI(title="Enterprise Intelligent Service Platform API", version="0.1.0")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])


@app.on_event("startup")
def startup() -> None:
    init_db()


def current_user(authorization: str | None = None) -> dict[str, Any]:
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(401, "请先登录")
    token = authorization.removeprefix("Bearer ").strip()
    token_hash = hashlib.sha256(token.encode("utf-8")).hexdigest()
    with closing(connect()) as conn:
        user = row(conn, "SELECT users.* FROM users JOIN auth_sessions ON auth_sessions.user_id = users.id WHERE auth_sessions.token_hash = ? AND auth_sessions.expires_at > ? AND users.status = 'active'", (token_hash, now()))
    if not user:
        raise HTTPException(401, "登录已失效，请重新登录")
    return public_user(user)

def get_user(authorization: str | None = Header(default=None)) -> dict[str, Any]:
    return current_user(authorization)

def require_role(user: dict[str, Any], *roles: str) -> None:
    if user["role"] not in roles:
        raise HTTPException(403, "permission denied")


@app.get("/api/health")
def health() -> dict[str, str]:
    return {"status": "ok", "service": "enterprise-platform-api", "version": "0.1.0"}


@app.post("/api/auth/login")
def login(payload: LoginRequest) -> dict[str, Any]:
    with closing(connect()) as conn:
        if payload.email.strip().lower() == "admin":
            user = row(conn, "SELECT * FROM users WHERE role = 'admin' AND status = 'active' ORDER BY id LIMIT 1")
        else:
            user = row(conn, "SELECT * FROM users WHERE lower(email) = lower(?) LIMIT 1", (payload.email.strip(),))
    if not user or user.get("status") != "active" or not verify_password(payload.password, user.get("password_hash", "")):
        raise HTTPException(401, "邮箱或密码错误")
    token = crypto_secrets.token_urlsafe(32)
    stamp = now()
    expires_at = datetime.fromtimestamp(datetime.now(timezone.utc).timestamp() + 8 * 3600, timezone.utc).isoformat(timespec="seconds")
    with closing(connect()) as conn:
        conn.execute("INSERT INTO auth_sessions (token_hash, user_id, created_at, expires_at) VALUES (?, ?, ?, ?)", (hashlib.sha256(token.encode("utf-8")).hexdigest(), user["id"], stamp, expires_at))
        conn.commit()
    return {"access_token": token, "token_type": "bearer", "user": public_user(user)}


@app.post("/api/auth/logout")
def logout(authorization: str | None = Header(default=None)) -> dict[str, bool]:
    if authorization and authorization.startswith("Bearer "):
        token = authorization.removeprefix("Bearer ").strip()
        with closing(connect()) as conn:
            conn.execute("DELETE FROM auth_sessions WHERE token_hash = ?", (hashlib.sha256(token.encode("utf-8")).hexdigest(),))
            conn.commit()
    return {"ok": True}

@app.post("/api/auth/forgot-password")
def forgot_password(payload: ForgotPasswordRequest) -> dict[str, Any]:
    """Create a review request without disclosing whether the account exists."""
    identifier = payload.identifier.strip()
    with closing(connect()) as conn:
        target = row(conn, "SELECT id, name, email, role, status FROM users WHERE lower(email) = lower(?) LIMIT 1", (identifier,))
        if target and target.get("role") != "admin" and target.get("status") == "active":
            pending = row(conn, "SELECT id FROM password_reset_requests WHERE user_id = ? AND status = 'pending' LIMIT 1", (target["id"],))
            if not pending:
                conn.execute(
                    "INSERT INTO password_reset_requests (id, user_id, identifier, status, created_at) VALUES (?, ?, ?, 'pending', ?)",
                    (f"PRQ-{uuid.uuid4().hex[:12].upper()}", target["id"], identifier, now()),
                )
                conn.commit()
    return {"ok": True, "message": "如果账号存在，管理员将收到密码重置通知。"}
@app.get("/api/auth/me")
def me(user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    return user


def visible_notifications(conn: sqlite3.Connection, user: dict[str, Any]) -> list[dict[str, Any]]:
    items = rows(conn, "SELECT id, title, detail, type, time, unread FROM notifications ORDER BY rowid DESC")
    if user.get('role') == 'admin':
        requests = rows(conn, "SELECT r.id, r.user_id, r.created_at, r.read_at, u.name, u.email FROM password_reset_requests r JOIN users u ON u.id = r.user_id WHERE r.status = 'pending' ORDER BY r.created_at DESC")
        for request in requests:
            items.insert(0, {'id': 'PR-' + request['id'], 'title': '员工密码重置申请', 'detail': f"{request['name']}（{request['email']}）请求重置登录密码", 'type': 'password_reset', 'time': request['created_at'], 'unread': request['read_at'] is None})
    for item in items:
        item['unread'] = bool(item['unread'])
    return items
@app.get("/api/bootstrap")
def bootstrap(user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    with closing(connect()) as conn:
        if user["role"] == "employee":
            leaves = rows(conn, "SELECT id, applicant, department, leave_type AS type, dates, days, status, reason, approver_id AS approverId, created_at AS createdAt, updated_at AS updatedAt FROM leave_requests WHERE applicant_id = ? ORDER BY created_at DESC", (user["id"],))
            documents = rows(conn, "SELECT id, name, owner, size, status, stage, updated, knowledge_base AS knowledgeBase FROM documents WHERE status = 'success' ORDER BY created_at DESC")
            tasks = []
            files = rows(conn, "SELECT id, name, file_type AS type, status, conversation_id AS conversationId, created_at AS createdAt, template, file_path AS filePath, owner_id FROM generated_files WHERE conversation_id IN (SELECT id FROM conversations WHERE user_id = ?) ORDER BY created_at DESC", (user["id"],))
        elif user["role"] == "manager":
            leaves = rows(conn, "SELECT id, applicant, department, leave_type AS type, dates, days, status, reason, approver_id AS approverId, created_at AS createdAt, updated_at AS updatedAt FROM leave_requests WHERE approver_id = ? ORDER BY created_at DESC", (user["id"],))
            documents = rows(conn, "SELECT id, name, owner, size, status, stage, updated, knowledge_base AS knowledgeBase FROM documents ORDER BY created_at DESC")
            tasks = rows(conn, "SELECT id, document, stage, status, progress, updated FROM ingestion_tasks ORDER BY updated DESC")
            files = rows(conn, "SELECT id, name, file_type AS type, status, conversation_id AS conversationId, created_at AS createdAt, template, file_path AS filePath, owner_id FROM generated_files ORDER BY created_at DESC")
        elif user["role"] == "manager":
            leaves = rows(conn, "SELECT id, applicant, department, leave_type AS type, dates, days, status, reason, approver_id AS approverId, created_at AS createdAt, updated_at AS updatedAt FROM leave_requests WHERE approver_id = ? ORDER BY created_at DESC", (user["id"],))
            documents = rows(conn, "SELECT id, name, owner, size, status, stage, updated, knowledge_base AS knowledgeBase FROM documents ORDER BY created_at DESC")
            tasks = rows(conn, "SELECT id, document, stage, status, progress, updated FROM ingestion_tasks ORDER BY updated DESC")
            files = rows(conn, "SELECT id, name, file_type AS type, status, conversation_id AS conversationId, created_at AS createdAt, template, file_path AS filePath, owner_id FROM generated_files ORDER BY created_at DESC")
        else:
            leaves = rows(conn, "SELECT id, applicant, department, leave_type AS type, dates, days, status, reason, approver_id AS approverId, created_at AS createdAt, updated_at AS updatedAt FROM leave_requests ORDER BY created_at DESC")
            documents = rows(conn, "SELECT id, name, owner, size, status, stage, updated, knowledge_base AS knowledgeBase FROM documents ORDER BY created_at DESC")
            tasks = rows(conn, "SELECT id, document, stage, status, progress, updated FROM ingestion_tasks ORDER BY updated DESC")
            files = rows(conn, "SELECT id, name, file_type AS type, status, conversation_id AS conversationId, created_at AS createdAt, template, file_path AS filePath, owner_id FROM generated_files ORDER BY created_at DESC")
        notifications = visible_notifications(conn, user)


        conversations = rows(conn, "SELECT id, title, preview, updated_at AS updated FROM conversations WHERE user_id = ? ORDER BY updated_at DESC", (user["id"],))
    return {"user": user, "leaves": leaves, "documents": documents, "tasks": tasks, "files": files, "notifications": notifications, "conversations": conversations}


@app.get("/api/team")
def list_managed_team(user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    require_role(user, "manager", "admin")
    with closing(connect()) as conn:
        if user["role"] == "admin":
            members = rows(conn, "SELECT u.id, u.name, u.email, u.department, u.role, u.status, u.supervisor_id AS supervisorId FROM users u WHERE u.id != ? AND (u.supervisor_id = ? OR u.department IN (SELECT name FROM departments WHERE manager_id = ?)) ORDER BY u.department, u.name", (user["id"], user["id"], user["id"]))
        else:
            members = rows(conn, "SELECT u.id, u.name, u.email, u.department, u.role, u.status, u.supervisor_id AS supervisorId FROM users u WHERE u.id != ? AND (u.supervisor_id = ? OR u.department IN (SELECT name FROM departments WHERE manager_id = ?)) ORDER BY u.department, u.name", (user["id"], user["id"], user["id"]))
    return {"members": members}
@app.get("/api/leaves")
def list_leaves(user: dict[str, Any] = Depends(get_user)) -> list[dict[str, Any]]:
    with closing(connect()) as conn:
        if user["role"] == "employee":
            return rows(conn, "SELECT id, applicant, department, leave_type AS type, dates, days, status, reason, approver_id AS approverId, created_at AS createdAt, updated_at AS updatedAt FROM leave_requests WHERE applicant_id = ? ORDER BY created_at DESC", (user["id"],))
        if user["role"] == "manager":
            return rows(conn, "SELECT id, applicant, department, leave_type AS type, dates, days, status, reason, approver_id AS approverId, created_at AS createdAt, updated_at AS updatedAt FROM leave_requests WHERE approver_id = ? ORDER BY created_at DESC", (user["id"],))
        if user["role"] == "manager":
            return rows(conn, "SELECT id, applicant, department, leave_type AS type, dates, days, status, reason, approver_id AS approverId, created_at AS createdAt, updated_at AS updatedAt FROM leave_requests WHERE approver_id = ? ORDER BY created_at DESC", (user["id"],))
        return rows(conn, "SELECT id, applicant, department, leave_type AS type, dates, days, status, reason, approver_id AS approverId, created_at AS createdAt, updated_at AS updatedAt FROM leave_requests ORDER BY created_at DESC")


@app.post("/api/leaves", status_code=201)
def create_leave(payload: LeaveCreate, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    leave_id = f"LV-{datetime.now().strftime('%Y%m%d')}-{uuid.uuid4().hex[:4].upper()}"
    stamp = now()

    with closing(connect()) as conn:
        try:
            start_date = datetime.strptime(payload.start, "%Y-%m-%d").date()
            end_date = datetime.strptime(payload.end, "%Y-%m-%d").date()
        except ValueError as exc:
            raise HTTPException(422, "日期格式必须为 YYYY-MM-DD") from exc
        if end_date < start_date:
            raise HTTPException(422, "结束日期不能早于开始日期")
        calculated_days = float((end_date - start_date).days + 1)
        approver_id = _resolve_approver_id(conn, user)
        conn.execute("INSERT INTO leave_requests VALUES (?, ?, ?, ?, ?, ?, ?, 'pending', ?, ?, ?, ?)", (leave_id, user["id"], user["name"], user["department"], payload.leave_type, f"{payload.start} 至 {payload.end}", calculated_days, payload.reason, approver_id, stamp, stamp))
        conn.execute("INSERT INTO notifications VALUES (?, ?, ?, 'approval', '刚刚', 1)", (f"N-{uuid.uuid4().hex[:8]}", "新的请假申请待审批", f"{user['name']} 提交了 {calculated_days:g} 天{payload.leave_type}申请"))
        audit(conn, user["id"], "leave.create", "leave_request", leave_id, payload.model_dump())
        conn.commit()
        created = row(conn, "SELECT id, applicant, department, leave_type AS type, dates, days, status, reason, approver_id AS approverId, created_at AS createdAt, updated_at AS updatedAt FROM leave_requests WHERE id = ?", (leave_id,))
    return created or {}


@app.post("/api/leaves/{leave_id}/decision")
def decide_leave(leave_id: str, payload: DecisionRequest, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    require_role(user, "manager", "admin")
    with closing(connect()) as conn:
        target = row(conn, "SELECT * FROM leave_requests WHERE id = ?", (leave_id,))
        if not target:
            raise HTTPException(404, "leave request not found")
        if user["role"] == "manager" and target.get("approver_id") != user["id"]:
            raise HTTPException(403, "该申请未分配给当前审批人")
        conn.execute("UPDATE leave_requests SET status = ?, updated_at = ? WHERE id = ?", (payload.action, now(), leave_id))
        conn.execute("INSERT INTO notifications VALUES (?, ?, ?, 'approval', '刚刚', 1)", (f"N-{uuid.uuid4().hex[:8]}", "请假审批结果已更新", f"申请单 {leave_id} 已被{payload.action}"))
        audit(conn, user["id"], f"leave.{payload.action}", "leave_request", leave_id, {"comment": payload.comment})
        conn.commit()
        return row(conn, "SELECT id, applicant, department, leave_type AS type, dates, days, status, reason, approver_id AS approverId, created_at AS createdAt, updated_at AS updatedAt FROM leave_requests WHERE id = ?", (leave_id,)) or {}


@app.get("/api/todos")
def list_todos(user: dict[str, Any] = Depends(get_user)) -> list[dict[str, Any]]:
    with closing(connect()) as conn:
        return rows(conn, "SELECT id, title, due_date AS dueDate, priority, status, created_at AS createdAt, updated_at AS updatedAt FROM todo_items WHERE owner_id = ? ORDER BY CASE status WHEN 'open' THEN 0 ELSE 1 END, CASE priority WHEN 'high' THEN 0 WHEN 'normal' THEN 1 ELSE 2 END, COALESCE(due_date, '9999-12-31'), created_at DESC", (user['id'],))

@app.post("/api/todos", status_code=201)
def create_todo(payload: TodoCreateRequest, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    due_date = payload.due_date.strip() if payload.due_date else None
    if due_date:
        try:
            datetime.strptime(due_date, "%Y-%m-%d")
        except ValueError as exc:
            raise HTTPException(422, "截止日期格式必须为 YYYY-MM-DD") from exc
    todo_id = f"TODO-{uuid.uuid4().hex[:10].upper()}"
    stamp = now()
    with closing(connect()) as conn:
        conn.execute("INSERT INTO todo_items VALUES (?, ?, ?, ?, ?, 'open', ?, ?)", (todo_id, user['id'], payload.title.strip(), due_date, payload.priority, stamp, stamp))
        audit(conn, user['id'], 'todo.create', 'todo', todo_id, payload.model_dump())
        conn.commit()
        return row(conn, "SELECT id, title, due_date AS dueDate, priority, status, created_at AS createdAt, updated_at AS updatedAt FROM todo_items WHERE id = ?", (todo_id,)) or {}

@app.patch("/api/todos/{todo_id}")
def update_todo(todo_id: str, payload: TodoUpdateRequest, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    with closing(connect()) as conn:
        target = row(conn, "SELECT id FROM todo_items WHERE id = ? AND owner_id = ?", (todo_id, user['id']))
        if not target:
            raise HTTPException(404, "待办不存在或无权操作")
        conn.execute("UPDATE todo_items SET status = ?, updated_at = ? WHERE id = ?", (payload.status, now(), todo_id))
        audit(conn, user['id'], 'todo.update', 'todo', todo_id, {'status': payload.status})
        conn.commit()
        return row(conn, "SELECT id, title, due_date AS dueDate, priority, status, created_at AS createdAt, updated_at AS updatedAt FROM todo_items WHERE id = ?", (todo_id,)) or {}

@app.delete("/api/todos/{todo_id}", status_code=204)
def delete_todo(todo_id: str, user: dict[str, Any] = Depends(get_user)) -> None:
    with closing(connect()) as conn:
        target = row(conn, "SELECT id FROM todo_items WHERE id = ? AND owner_id = ?", (todo_id, user['id']))
        if not target:
            raise HTTPException(404, "待办不存在或无权操作")
        conn.execute("DELETE FROM todo_items WHERE id = ?", (todo_id,))
        audit(conn, user['id'], 'todo.delete', 'todo', todo_id, None)
        conn.commit()

@app.get("/api/calendar")
def calendar_items(month: str, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    try:
        start = datetime.strptime(month + "-01", "%Y-%m-%d").date()
    except ValueError as exc:
        raise HTTPException(422, "月份格式必须为 YYYY-MM") from exc
    end = (start.replace(day=28) + timedelta(days=4)).replace(day=1) - timedelta(days=1)
    with closing(connect()) as conn:
        todos = rows(conn, "SELECT id, title, due_date AS dueDate, priority, status FROM todo_items WHERE owner_id = ? AND due_date >= ? AND due_date <= ?", (user['id'], start.isoformat(), end.isoformat()))
        if user['role'] == 'employee':
            leaves = rows(conn, "SELECT id, applicant, leave_type AS type, dates, status FROM leave_requests WHERE applicant_id = ?", (user['id'],))
        else:
            leaves = rows(conn, "SELECT id, applicant, leave_type AS type, dates, status FROM leave_requests WHERE applicant_id = ? OR approver_id = ?", (user['id'], user['id']))
    events: list[dict[str, Any]] = []
    for todo in todos:
        events.append({'id': todo['id'], 'kind': 'todo', 'title': todo['title'], 'date': todo['dueDate'], 'priority': todo['priority'], 'status': todo['status']})
    for leave in leaves:
        parts = str(leave.get('dates') or '').split(' 至 ')
        if len(parts) != 2:
            continue
        try:
            leave_start = datetime.strptime(parts[0], "%Y-%m-%d").date()
            leave_end = datetime.strptime(parts[1], "%Y-%m-%d").date()
        except ValueError:
            continue
        if leave_end < start or leave_start > end:
            continue
        events.append({'id': leave['id'], 'kind': 'leave', 'title': f"{leave['applicant']} · {leave['type']}", 'start': parts[0], 'end': parts[1], 'status': leave['status']})
    return {'month': month, 'events': events}
@app.get("/api/knowledge/documents")
def list_documents(user: dict[str, Any] = Depends(get_user)) -> list[dict[str, Any]]:
    with closing(connect()) as conn:
        if user["role"] == "employee":
            return rows(conn, "SELECT id, name, owner, size, status, stage, updated, knowledge_base AS knowledgeBase FROM documents WHERE status = 'success' ORDER BY created_at DESC")
        return rows(conn, "SELECT id, name, owner, size, status, stage, updated, knowledge_base AS knowledgeBase FROM documents ORDER BY created_at DESC")


@app.post("/api/knowledge/documents", status_code=201)
async def upload_document(file: UploadFile = File(...), user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    require_role(user, "admin", "manager")
    document_id = f"DOC-{uuid.uuid4().hex[:8].upper()}"
    task_id = f"ING-{uuid.uuid4().hex[:8].upper()}"
    size = "未知"
    content = await file.read()
    if content:
        size = f"{len(content) / 1024 / 1024:.1f} MB"
    stamp = now()

    with closing(connect()) as conn:
        conn.execute("INSERT INTO documents VALUES (?, ?, ?, ?, 'processing', '上传中', '刚刚', '研发知识库', ?)", (document_id, file.filename or "未命名文件", user["name"], size, stamp))
        conn.execute("INSERT INTO ingestion_tasks VALUES (?, ?, ?, '上传', 'processing', 12, '刚刚')", (task_id, document_id, file.filename or "未命名文件"))
        audit(conn, user["id"], "document.upload", "document", document_id, {"name": file.filename, "task_id": task_id})
        conn.commit()
    return {"id": document_id, "task_id": task_id, "name": file.filename, "status": "processing", "stage": "上传中"}


@app.get("/api/knowledge/tasks")
def list_tasks(user: dict[str, Any] = Depends(get_user)) -> list[dict[str, Any]]:
    with closing(connect()) as conn:
        return rows(conn, "SELECT id, document, stage, status, progress, updated FROM ingestion_tasks ORDER BY updated DESC")


@app.get("/api/notifications")
def list_notifications(user: dict[str, Any] = Depends(get_user)) -> list[dict[str, Any]]:
    with closing(connect()) as conn:
        return visible_notifications(conn, user)


@app.post("/api/notifications/{notification_id}/read")
def mark_notification_read(notification_id: str, user: dict[str, Any] = Depends(get_user)) -> dict[str, bool]:
    with closing(connect()) as conn:
        if notification_id.startswith("PR-"):
            require_role(user, "admin")
            request_id = notification_id.removeprefix("PR-")
            conn.execute("UPDATE password_reset_requests SET read_at = COALESCE(read_at, ?) WHERE id = ? AND status = 'pending'", (now(), request_id))
        else:
            conn.execute("UPDATE notifications SET unread = 0 WHERE id = ?", (notification_id,))
        audit(conn, user["id"], "notification.read", "notification", notification_id)
        conn.commit()
    return {"ok": True}

def ensure_conversation(conn: sqlite3.Connection, conversation_id: str, user_id: str, title: str = "新对话") -> None:
    if not row(conn, "SELECT id FROM conversations WHERE id = ?", (conversation_id,)):
        stamp = now()
        conn.execute("INSERT INTO conversations VALUES (?, ?, ?, '', 'active', ?, ?)", (conversation_id, user_id, title, stamp, stamp))


def resolve_model_config(conn: sqlite3.Connection) -> dict[str, Any]:
    config = row(conn, "SELECT * FROM model_configs WHERE scope_type = 'global' AND scope_id = 'platform' LIMIT 1")
    if not config or not config.get("api_key_ciphertext"):
        raise HTTPException(409, "模型配置未保存可用的 API Key，请重新保存模型配置")
    try:
        api_key = decrypt_secret(config["api_key_ciphertext"])
    except ValueError as exc:
        raise HTTPException(409, "模型 API Key 无法解密，请重新保存模型配置") from exc
    if not config.get("api_url") or not config.get("model"):
        raise ValueError("模型配置缺少 API URL 或模型名称")
    config["api_key"] = api_key
    return config


def model_chat_endpoint(api_url: str, api_format: str) -> str:
    base = api_url.strip().rstrip("/")
    if base.endswith(("/chat/completions", "/responses", "/messages")):
        return base
    if api_format == "anthropic-messages":
        return f"{base}/messages"
    if api_format == "openai-responses":
        return f"{base}/responses"
    return f"{base}/chat/completions"


def extract_model_response(payload: Any, api_format: str) -> str:
    if api_format == "anthropic-messages":
        content = payload.get("content") if isinstance(payload, dict) else None
        if isinstance(content, list):
            return "".join(str(item.get("text", "")) for item in content if isinstance(item, dict)).strip()
    if api_format == "openai-responses":
        output_text = payload.get("output_text") if isinstance(payload, dict) else None
        if isinstance(output_text, str) and output_text.strip():
            return output_text.strip()
        output = payload.get("output") if isinstance(payload, dict) else None
        if isinstance(output, list):
            chunks: list[str] = []
            for item in output:
                for content in item.get("content", []) if isinstance(item, dict) else []:
                    if isinstance(content, dict) and isinstance(content.get("text"), str):
                        chunks.append(content["text"])
            return "".join(chunks).strip()
    choices = payload.get("choices") if isinstance(payload, dict) else None
    if isinstance(choices, list) and choices:
        message = choices[0].get("message", {})
        content = message.get("content") if isinstance(message, dict) else None
        if isinstance(content, str):
            return content.strip()
        if isinstance(content, list):
            return "".join(str(item.get("text", "")) for item in content if isinstance(item, dict)).strip()
    return ""


def extract_model_usage(payload: Any) -> dict[str, int]:
    usage = payload.get("usage") if isinstance(payload, dict) else None
    if not isinstance(usage, dict):
        return {"input_tokens": 0, "output_tokens": 0, "total_tokens": 0}
    input_tokens = usage.get("prompt_tokens", usage.get("input_tokens", 0))
    output_tokens = usage.get("completion_tokens", usage.get("output_tokens", 0))
    total_tokens = usage.get("total_tokens", 0)
    try:
        input_value = max(0, int(input_tokens or 0))
        output_value = max(0, int(output_tokens or 0))
        total_value = max(0, int(total_tokens or 0)) or input_value + output_value
    except (TypeError, ValueError):
        return {"input_tokens": 0, "output_tokens": 0, "total_tokens": 0}
    return {"input_tokens": input_value, "output_tokens": output_value, "total_tokens": total_value}


def call_configured_model(config: dict[str, Any], messages: list[dict[str, str]]) -> tuple[str, int, str, dict[str, int]]:
    api_format = config.get("api_format") or "openai-responses"
    if api_format == "google-generative-ai":
        raise ValueError("Google Generative AI 对话适配尚未启用，请选择 OpenAI 或 Anthropic 兼容协议")
    endpoint = model_chat_endpoint(config["api_url"], api_format)
    if api_format == "anthropic-messages":
        headers = {"Accept": "application/json", "Content-Type": "application/json", "x-api-key": config["api_key"], "anthropic-version": "2023-06-01"}
        body = {"model": config["model"], "max_tokens": 2048, "messages": [item for item in messages if item["role"] != "system"]}
    elif api_format == "openai-responses":
        headers = {"Accept": "application/json", "Content-Type": "application/json", "Authorization": f"Bearer {config['api_key']}"}
        body = {"model": config["model"], "input": messages, "stream": False}
    else:
        headers = {"Accept": "application/json", "Content-Type": "application/json", "Authorization": f"Bearer {config['api_key']}"}
        body = {"model": config["model"], "messages": messages, "stream": False}
    request = urllib.request.Request(endpoint, headers=headers, data=json.dumps(body, ensure_ascii=False).encode("utf-8"), method="POST")
    started = time.perf_counter()
    try:
        with urllib.request.urlopen(request, timeout=90) as response:
            raw = response.read(8 * 1024 * 1024)
    except urllib.error.HTTPError as exc:
        raw = exc.read(4096).decode("utf-8", errors="replace")
        safe = raw.replace(config["api_key"], "[REDACTED]")
        raise ValueError(f"模型接口 HTTP {exc.code}: {safe[:512]}") from exc
    except urllib.error.URLError as exc:
        raise ValueError(f"模型接口连接失败: {exc.reason}") from exc
    except TimeoutError as exc:
        raise ValueError("模型接口请求超时") from exc
    try:
        payload = json.loads(raw.decode("utf-8"))
    except (json.JSONDecodeError, UnicodeDecodeError) as exc:
        raise ValueError("模型接口返回内容不是有效 JSON") from exc
    response_text = extract_model_response(payload, api_format)
    if not response_text:
        raise ValueError("模型接口返回成功，但没有找到可显示的文本内容")
    return response_text, int((time.perf_counter() - started) * 1000), endpoint, extract_model_usage(payload)


def _extract_leave_details(request_text: str) -> dict[str, str]:
    text = request_text
    leave_types = ("年假", "病假", "事假", "调休", "婚假", "产假", "陪产假", "育儿假")
    leave_type = next((item for item in leave_types if item in text), "请假")
    date_pattern = r"(20\d{2})[./-](\d{1,2})[./-](\d{1,2})\s*(?:至|到|--|-)\s*(?:(20\d{2})[./-])?(\d{1,2})[./-](\d{1,2})"
    date_match = re.search(date_pattern, text)
    labeled_dates = re.findall(r"20\d{2}[./-]\d{1,2}[./-]\d{1,2}", text)
    days_match = re.search(r"(?:连续|共|请假)\s*(\d+(?:\.\d+)?)\s*天", text)
    days = float(days_match.group(1)) if days_match else None
    if date_match:
        year, month, day, end_year, end_month, end_day = date_match.groups()
        start_date = datetime(int(year), int(month), int(day)).date()
        end_date = datetime(int(end_year or year), int(end_month), int(end_day)).date()
        days = days or float((end_date - start_date).days + 1)
    elif len(labeled_dates) >= 2:
        start_date = datetime.strptime(labeled_dates[0].replace('/', '-').replace('.', '-'), '%Y-%m-%d').date()
        end_date = datetime.strptime(labeled_dates[1].replace('/', '-').replace('.', '-'), '%Y-%m-%d').date()
        days = days or float((end_date - start_date).days + 1)
    else:
        relative = re.search(r"(明天|后天)(?:起|开始)?(?:连续|共)?\s*(\d+(?:\.\d+)?)?\s*天", text)
        if relative:
            offset = 1 if relative.group(1) == "明天" else 2
            days = days or float(relative.group(2) or 1)
            start_date = datetime.now().date() + timedelta(days=offset)
            end_date = start_date + timedelta(days=int(days) - 1)
        else:
            start_date = None
            end_date = None
    handover_match = re.search(r"((?:已|已经|并已|并且已|我已)?安排(?:好)?[^。；\n]{0,24}?工作交接[^。；\n]{0,24})", text)
    handover = handover_match.group(1).strip(" ，,；;") if handover_match else "未填写"
    reason_match = re.search(r"(?:请假理由|申请事由|事由|原因)[:：]?\s*([^。；\n]+)", text)
    reason = reason_match.group(1).strip() if reason_match else ""
    if not reason:
        reason_match = re.search(r"(?:因|由于)([^。；\n]+)", text)
        reason = reason_match.group(1).strip() if reason_match else "未填写"
    reason = re.sub(r"[，,；;]\s*(?:已|已经|并已|并且已|我已)?安排(?:好)?[^。；\n]*工作交接[^。；\n]*", "", reason).strip()
    reason = re.sub(r"[，,；;]\s*(?:并)?(?:生成|制作|导出|创建)[^。；\n]*$", "", reason).strip()
    reason = re.sub(r"^(?:是|为)\s*", "", reason).strip()
    if not reason:
        reason = "未填写"
    return {
        "leave_type": leave_type,
        "start": start_date.isoformat() if start_date else "未填写",
        "end": end_date.isoformat() if end_date else "未填写",
        "days": f"{days:g}" if days is not None else "未填写",
        "reason": reason,
        "handover": handover,
    }


def _resolve_approver_id(conn: sqlite3.Connection, user: dict[str, Any]) -> str:
    supervisor_id = user.get("supervisor_id")
    if supervisor_id:
        supervisor = row(conn, "SELECT id FROM users WHERE id = ? AND status = 'active'", (supervisor_id,))
        if supervisor:
            return supervisor["id"]
    department = row(conn, "SELECT manager_id FROM departments WHERE name = ?", (user["department"],))
    manager_id = (department or {}).get("manager_id")
    if manager_id and manager_id != user["id"]:
        return manager_id
    raise HTTPException(422, "当前用户未配置直属审批人")

def _ensure_leave_request(conn: sqlite3.Connection, user: dict[str, Any], details: dict[str, str]) -> tuple[str, str]:
    if details["start"] == "未填写" or details["end"] == "未填写":
        raise HTTPException(422, "请假申请缺少起止日期，无法提交审批")
    try:
        start_date = datetime.strptime(details["start"], "%Y-%m-%d").date()
        end_date = datetime.strptime(details["end"], "%Y-%m-%d").date()
    except ValueError as exc:
        raise HTTPException(422, "请假日期格式无效") from exc
    if end_date < start_date:
        raise HTTPException(422, "结束日期不能早于开始日期")
    details["days"] = f"{(end_date - start_date).days + 1:g}"
    approver_id = _resolve_approver_id(conn, user)
    dates = f"{details['start']} 至 {details['end']}"
    existing = row(conn, "SELECT id FROM leave_requests WHERE applicant_id = ? AND leave_type = ? AND dates = ? AND reason = ? AND status = 'pending' ORDER BY created_at DESC LIMIT 1", (user["id"], details["leave_type"], dates, details["reason"]))
    if existing:
        return existing["id"], approver_id
    leave_id = f"LV-{datetime.now().strftime('%Y%m%d')}-{uuid.uuid4().hex[:6].upper()}"
    stamp = now()
    conn.execute("INSERT INTO leave_requests VALUES (?, ?, ?, ?, ?, ?, ?, 'pending', ?, ?, ?, ?)", (leave_id, user["id"], user["name"], user["department"], details["leave_type"], dates, float(details["days"]), details["reason"], approver_id, stamp, stamp))
    conn.execute("INSERT INTO notifications VALUES (?, ?, ?, 'approval', '刚刚', 1)", (f"N-{uuid.uuid4().hex[:8]}", "新的请假申请待审批", f"{user['name']} 提交了 {details['days']} 天{details['leave_type']}申请"))
    audit(conn, user["id"], "leave.create_from_document", "leave_request", leave_id, details)
    return leave_id, approver_id

def _create_docx_file(conversation_id: str, user: dict[str, Any], request_text: str, template: str = "企业服务申请", leave_request_id: str | None = None) -> tuple[str, str]:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    output_dir = DATA_DIR / "generated"
    output_dir.mkdir(parents=True, exist_ok=True)
    file_id = f"FILE-{uuid.uuid4().hex[:10].upper()}"
    is_leave = any(keyword in request_text for keyword in ("请假", "病假", "年假", "事假", "调休"))
    title = "请假申请单" if is_leave else template
    filename = f"{user['name']}-{title}-{file_id}.docx"
    path = output_dir / filename
    document = Document()
    heading = document.add_heading(title, 0)
    heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta = document.add_paragraph()
    meta.add_run("生成时间：").bold = True
    meta.add_run(now())
    meta.add_run("    ")
    meta.add_run("conversation_id：").bold = True
    meta.add_run(conversation_id)
    if is_leave:
        details = _extract_leave_details(request_text)
        table = document.add_table(rows=0, cols=2)
        table.style = "Table Grid"
        fields = [
            ("申请单号", leave_request_id or "未关联"),
            ("申请人", user["name"]),
            ("所属部门", user["department"]),
            ("请假类型", details["leave_type"]),
            ("开始日期", details["start"]),
            ("结束日期", details["end"]),
            ("请假天数", details["days"] + " 天" if details["days"] != "未填写" else "未填写"),
            ("申请事由", details["reason"]),
            ("工作交接", details["handover"]),
            ("审批状态", "待直属主管审批"),
            ("审批人", "由组织关系自动匹配"),
        ]
        for label, value in fields:
            cells = table.add_row().cells
            cells[0].text = label
            cells[1].text = value
        document.add_paragraph("说明：本申请单由企业智能服务平台根据当前会话生成，审批结果以审批中心记录为准。")
    else:
        document.add_heading("申请内容", level=1)
        document.add_paragraph(request_text or "未填写")
        document.add_heading("处理说明", level=1)
        document.add_paragraph("本文件由企业智能服务平台根据当前会话生成，具体业务审批结果以业务系统记录为准。")
    document.save(path)
    return file_id, str(path)

def _insert_agent_event(conn: sqlite3.Connection, conversation_id: str, task_id: str, trace_id: str, seq: int, event_type: str, label: str, status: str, duration: str, summary: str, data: dict[str, Any]) -> None:
    conn.execute("INSERT INTO agent_events VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)", (f"EVT-{uuid.uuid4().hex[:10].upper()}", conversation_id, task_id, trace_id, seq, event_type, label, status, duration, summary, json.dumps(data, ensure_ascii=False), now()))


def _email_request_kind(text: str) -> str | None:
    lowered = text.lower()
    if any(keyword in lowered for keyword in ('读邮件', '读取邮件', '查邮件', '查看邮件', '收邮件', '收件箱', '最近邮件', '最近收到', '最新邮件', '最新一封', '最新的', '最后一封', '新邮件', '邮件列表', '未读邮件', '未读邮箱')) or ('邮件' in lowered and any(keyword in lowered for keyword in ('读取', '查看', '查询', '最近', '最新', '收取', '未读'))):
        return 'read'
    if any(keyword in lowered for keyword in ('发邮件', '发送邮件', '写邮件', '发一封邮件', '发信')):
        return 'send'
    return None


def _parse_email_send_request(text: str) -> tuple[list[str], str, str]:
    recipients = list(dict.fromkeys(re.findall(r'[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}', text)))
    subject_match = re.search(r'(?:主题|标题)\s*[:：]\s*(.*?)(?=\s*(?:正文|内容)\s*[:：]|$)', text, re.S | re.I)
    body_match = re.search(r'(?:正文|内容)\s*[:：]\s*(.+)$', text, re.S | re.I)
    subject = subject_match.group(1).strip() if subject_match else ''
    body = body_match.group(1).strip() if body_match else ''
    return recipients, subject, body


def _conversation_task_events(conversation_id: str, task_id: str) -> list[dict[str, Any]]:
    with closing(connect()) as conn:
        return rows(conn, "SELECT id, seq, label, event_type AS type, status, duration, safe_summary AS summary, data FROM agent_events WHERE conversation_id = ? AND task_id = ? ORDER BY seq", (conversation_id, task_id))
@app.get("/api/conversations")
def list_conversations(user: dict[str, Any] = Depends(get_user)) -> list[dict[str, Any]]:
    with closing(connect()) as conn:
        return rows(conn, "SELECT id, title, preview, updated_at AS updated FROM conversations WHERE user_id = ? ORDER BY updated_at DESC", (user["id"],))


@app.get("/api/conversations/{conversation_id}")
def conversation_detail(conversation_id: str, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    with closing(connect()) as conn:
        conversation = row(conn, "SELECT id, user_id, title, preview, status, created_at, updated_at AS updated FROM conversations WHERE id = ?", (conversation_id,))
        if not conversation:
            raise HTTPException(404, "conversation not found")
        if conversation["user_id"] != user["id"] and user["role"] not in {"manager", "admin"}:
            raise HTTPException(403, "conversation access denied")
        messages = rows(conn, "SELECT id, role, content, seq, created_at FROM messages WHERE conversation_id = ? ORDER BY seq", (conversation_id,))
        attachment_rows = rows(conn, "SELECT attachment_id AS id, filename AS name, size, content_type AS contentType, provider_message_id AS messageId, account_id AS accountId, message_seq FROM conversation_attachments WHERE conversation_id = ? ORDER BY message_seq, id", (conversation_id,))
        attachments_by_seq: dict[int, list[dict[str, Any]]] = {}
        for attachment in attachment_rows:
            attachments_by_seq.setdefault(int(attachment.get('message_seq', 0)), []).append(attachment)
        for message in messages:
            message['attachments'] = attachments_by_seq.get(int(message.get('seq', 0)), [])
        events = rows(conn, "SELECT id, task_id, trace_id, seq, label, event_type AS type, status, duration, safe_summary AS summary, data, created_at FROM agent_events WHERE conversation_id = ? ORDER BY seq", (conversation_id,))
        audit_items = rows(conn, "SELECT id, action, resource_type, resource_id, detail, created_at FROM audit_events WHERE resource_id = ? ORDER BY created_at", (conversation_id,)) if user["role"] in {"manager", "admin"} else []
    return {"conversation": conversation, "messages": messages, "events": events, "audit": audit_items}

@app.post("/api/conversations/{conversation_id}/messages")
def send_message(conversation_id: str, payload: MessageRequest, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    task_id = f"TASK-{uuid.uuid4().hex[:10].upper()}"
    trace_id = f"TR-{uuid.uuid4().hex[:8].upper()}"
    stamp = now()
    request_lower = payload.content.lower()
    with closing(connect()) as conn:
        existing = row(conn, "SELECT id, user_id FROM conversations WHERE id = ?", (conversation_id,))
        if existing and existing["user_id"] != user["id"] and user["role"] not in {"manager", "admin"}:
            raise HTTPException(403, "conversation access denied")
        ensure_conversation(conn, conversation_id, user["id"], payload.content[:28])
        message_count = int(row(conn, "SELECT COALESCE(MAX(seq), 0) AS value FROM messages WHERE conversation_id = ?", (conversation_id,))["value"]) + 1
        event_seq = int(row(conn, "SELECT COALESCE(MAX(seq), 0) AS value FROM agent_events WHERE conversation_id = ?", (conversation_id,))["value"]) + 1
        conn.execute("INSERT INTO messages VALUES (?, ?, 'user', ?, ?, ?)", (f"MSG-{uuid.uuid4().hex[:10].upper()}", conversation_id, payload.content, message_count, stamp))
        _insert_agent_event(conn, conversation_id, task_id, trace_id, event_seq, "intent_detected", "识别用户请求", "success", "14ms", "识别为企业服务请求", {"intent": "enterprise_service"})
        _insert_agent_event(conn, conversation_id, task_id, trace_id, event_seq + 1, "permission_checked", "权限校验", "success", "18ms", f"允许 {user['role']} 继续执行", {"role": user["role"], "allowed": True})
        history = rows(conn, "SELECT role, content FROM messages WHERE conversation_id = ? ORDER BY seq DESC LIMIT 12", (conversation_id,))[::-1]
        history_text = chr(10).join(str(item["content"]) for item in history)
        history_lower = history_text.lower()
        continuation = any(keyword in request_lower for keyword in ("继续", "接着", "完成", "写好", "做成", "生成", "提交"))
        document_context = any(keyword in history_lower for keyword in ("docx", "word", "请假条", "文件", "请假申请", "生成 office 文档", "生成一个文档"))
        document_action = any(keyword in request_lower for keyword in ("生成", "写好", "制作", "导出", "创建"))
        document_request = (document_action and any(keyword in request_lower for keyword in ("文档", "docx", "word", "文件", "请假条"))) or (continuation and document_context and any(keyword in request_lower for keyword in ("请假", "病假", "年假", "文档", "文件", "它", "这个", "提交")))
        config = None if document_request else resolve_model_config(conn)
        conn.commit()
    email_kind = _email_request_kind(payload.content)
    if not email_kind and '邮件' in history_lower and (re.search(r'(?:第\s*)?(?:[一二两三四五六七八九十百千万]+|\d+)\s*(?:封|条|个)', payload.content) or any(keyword in request_lower for keyword in ('正文', '全文', '邮件内容', '第一条内容', '第二条内容', '你再看看', '再看看', '再看一下', '具体看看', '打开附件'))):
        email_kind = 'read'
    confirm_email = any(keyword in request_lower for keyword in ('确认', '确认发送', '好的发', '可以发', '发吧')) and not any(keyword in request_lower for keyword in ('取消', '不要发', '不发'))
    pending_email: dict[str, Any] | None = None
    if confirm_email:
        with closing(connect()) as conn:
            pending_event = row(conn, "SELECT data FROM agent_events WHERE conversation_id = ? AND event_type = 'email.send_pending' AND status = 'pending' ORDER BY seq DESC LIMIT 1", (conversation_id,))
        if pending_event:
            try:
                candidate = json.loads(pending_event.get('data') or '{}')
                pending_email = candidate if isinstance(candidate, dict) else None
            except json.JSONDecodeError:
                pending_email = None

    if pending_email and confirm_email:
        email_payload = AgentMailSendRequest(to=[str(item) for item in pending_email.get('to', [])], subject=str(pending_email.get('subject', '')), body=str(pending_email.get('body', '')), confirmed=True, account_id=pending_email.get('account_id'))
        try:
            result = _send_agentmail_message(email_payload, user, 'agent')
            response_text = '邮件已发送成功，已保存到企业邮箱的已发送记录。收件人：' + '、'.join(email_payload.to)
            with closing(connect()) as conn:
                _insert_agent_event(conn, conversation_id, task_id, trace_id, event_seq + 2, 'tool.email_send', '通过 AgentMail 发送邮件', 'success', '已完成', response_text, {'to': email_payload.to, 'subject': email_payload.subject, 'message_id': result.get('providerId', '')})
                conn.execute("UPDATE agent_events SET status = 'completed' WHERE conversation_id = ? AND event_type = 'email.send_pending' AND status = 'pending'", (conversation_id,))
                conn.execute("INSERT INTO messages VALUES (?, ?, 'assistant', ?, ?, ?)", (f'MSG-{uuid.uuid4().hex[:10].upper()}', conversation_id, response_text, message_count + 1, now()))
                conn.execute("UPDATE conversations SET preview = ?, updated_at = ? WHERE id = ?", (payload.content[:60], now(), conversation_id))
                audit(conn, user['id'], 'email.send', 'conversation', conversation_id, {'task_id': task_id, 'to': email_payload.to, 'subject': email_payload.subject, 'source': 'agent'})
                conn.commit()
            return {'conversation_id': conversation_id, 'task_id': task_id, 'trace_id': trace_id, 'response': response_text, 'events': _conversation_task_events(conversation_id, task_id)}
        except (RuntimeError, HTTPException) as exc:
            with closing(connect()) as conn:
                _insert_agent_event(conn, conversation_id, task_id, trace_id, event_seq + 2, 'tool.email_send', '通过 AgentMail 发送邮件', 'failed', '--', str(exc), {'error': str(exc)})
                audit(conn, user['id'], 'email.send_failed', 'conversation', conversation_id, {'task_id': task_id, 'error': str(exc)})
                conn.commit()
            raise HTTPException(status_code=502, detail=f'邮件发送失败：{exc}') from exc

    if email_kind == 'read':
        try:
            account_id = _active_account_id(user)
            requested_match = re.search(r'(?:第\s*)?([一二两三四五六七八九十百千万]+|\d+)\s*(?:封|条|个)', payload.content)
            number_words = {'一': 1, '二': 2, '两': 2, '三': 3, '四': 4, '五': 5, '六': 6, '七': 7, '八': 8, '九': 9, '十': 10}
            requested_index = None
            if requested_match:
                token = requested_match.group(1)
                requested_index = int(token) if token.isdigit() else number_words.get(token)
            latest_requested = any(keyword in request_lower for keyword in ('最新', '最后一封', '你再看看', '再看看', '再看一下', '具体看看'))
            detail_requested = requested_index is not None or latest_requested or any(keyword in request_lower for keyword in ('正文', '全文', '邮件内容', '内容', '附件'))
            context_ids: list[str] = []
            context_unread = False
            if detail_requested:
                with closing(connect()) as conn:
                    context_event = row(conn, "SELECT data FROM agent_events WHERE conversation_id = ? AND event_type = 'tool.email_list' ORDER BY seq DESC LIMIT 1", (conversation_id,))
                if context_event:
                    try:
                        context_data = json.loads(context_event.get('data') or '{}')
                        context_ids = [str(item) for item in context_data.get('message_ids', []) if item]
                        context_unread = bool(context_data.get('unread_only'))
                    except (TypeError, json.JSONDecodeError):
                        context_ids = []
            unread_only = '未读' in request_lower or (detail_requested and context_unread)
            detail_message = None
            email_attachments: list[dict[str, Any]] = []
            provider_items: list[dict[str, Any]] = []
            remote_id = ''
            if latest_requested and not requested_index:
                requested_index = 1
            if detail_requested and not latest_requested and context_ids and requested_index and requested_index <= len(context_ids):
                remote_id = context_ids[requested_index - 1]
            else:
                list_arguments = ['message', '+list', '--dir', 'inbox', '--limit', '8']
                if unread_only:
                    list_arguments.append('--is-unread')
                data = _run_agentmail(list_arguments, workspace=_account_workspace(user))
                provider_items = data.get('data') if isinstance(data.get('data'), list) else []
                if detail_requested and requested_index and requested_index <= len(provider_items):
                    selected_item = provider_items[requested_index - 1] if isinstance(provider_items[requested_index - 1], dict) else {}
                    remote_id = str(selected_item.get('message_id') or selected_item.get('id') or '')
            if detail_requested and remote_id:
                detail_data = _run_agentmail(['message', '+read', '--id', remote_id], workspace=_account_workspace(user))
                detail_message = detail_data.get('message') if isinstance(detail_data.get('message'), dict) else detail_data
                if isinstance(detail_message, dict):
                    provider_items = [detail_message]
                    for attachment in detail_message.get('attachments') or []:
                        if not isinstance(attachment, dict):
                            continue
                        attachment_id = str(attachment.get('attachment_id') or attachment.get('id') or '')
                        filename = str(attachment.get('filename') or attachment.get('name') or '附件')
                        item = {'id': attachment_id, 'name': filename, 'size': attachment.get('size') or 0, 'contentType': str(attachment.get('content_type') or attachment.get('mime_type') or 'application/octet-stream'), 'messageId': remote_id, 'accountId': account_id}
                        if attachment.get('download_url'):
                            item['downloadUrl'] = str(attachment['download_url'])
                        if attachment_id or item.get('downloadUrl'):
                            email_attachments.append(item)
            with closing(connect()) as conn:
                for item in provider_items:
                    if isinstance(item, dict):
                        _upsert_email(conn, user['id'], item, 'inbox', 'agent', account_id=account_id)
                message_ids = [str(item.get('message_id') or item.get('id') or '') for item in provider_items if isinstance(item, dict)]
                if detail_message and isinstance(detail_message, dict):
                    sender = detail_message.get('from')
                    sender_email = sender.get('email') if isinstance(sender, dict) else sender
                    raw_body = str(detail_message.get('body') or detail_message.get('text') or detail_message.get('snippet') or '')
                    decoded_body = html_lib.unescape(raw_body).replace('\\u003c', '<').replace('\\u003e', '>')
                    decoded_body = re.sub(r'<style[^>]*>.*?</style>', ' ', decoded_body, flags=re.S | re.I)
                    plain_body = re.sub(r'<[^>]+>', ' ', decoded_body)
                    plain_body = re.sub(r'\.qmbox\b.*?(?=Agent Mail 接入成功|邮箱地址)', ' ', plain_body, flags=re.S | re.I)
                    plain_body = re.sub(r'---\s*原始邮件\s*---.*$', '', plain_body, flags=re.S | re.I)
                    plain_body = re.sub(r'\s+', ' ', plain_body).strip()
                    ordinal = requested_index or 1
                    title = '最新未读邮件' if latest_requested and unread_only else '最新一封邮件' if latest_requested else '第' + str(ordinal) + ('封未读邮件' if unread_only else '封邮件')
                    response_text = title + '：' + '\n主题：' + str(detail_message.get('subject') or '(无主题)') + '\n发件人：' + str(sender_email or '未知发件人') + '\n正文：\n' + (plain_body[:4000] or '（正文为空）')
                    _insert_agent_event(conn, conversation_id, task_id, trace_id, event_seq + 2, 'tool.email_read', '读取邮件正文', 'success', '已完成', f'已读取第 {ordinal} 封邮件正文', {'message_id': remote_id, 'index': ordinal, 'unread_only': unread_only, 'attachments': email_attachments})
                elif provider_items:
                    lines = [f"{index + 1}. {str(item.get('subject') or '(无主题)')} · {str((item.get('from') or {}).get('email') if isinstance(item.get('from'), dict) else item.get('from') or '未知发件人')}" for index, item in enumerate(provider_items[:8])]
                    response_text = ('最近未读邮件：' if unread_only else '最近收到的邮件：') + '\n' + '\n'.join(lines)
                    _insert_agent_event(conn, conversation_id, task_id, trace_id, event_seq + 2, 'tool.email_list', '读取企业邮箱', 'success', '已完成', f'读取到 {len(provider_items)} 封' + ('未读' if unread_only else '') + '邮件', {'folder': 'inbox', 'count': len(provider_items), 'message_ids': message_ids, 'unread_only': unread_only, 'account_id': account_id})
                else:
                    response_text = ('当前收件箱没有未读邮件。' if unread_only else '当前收件箱没有邮件。')
                    _insert_agent_event(conn, conversation_id, task_id, trace_id, event_seq + 2, 'tool.email_list', '读取企业邮箱', 'success', '已完成', response_text, {'folder': 'inbox', 'count': 0, 'message_ids': [], 'unread_only': unread_only, 'account_id': account_id})
                conn.execute("INSERT INTO messages VALUES (?, ?, 'assistant', ?, ?, ?)", (f'MSG-{uuid.uuid4().hex[:10].upper()}', conversation_id, response_text, message_count + 1, now()))
                conn.execute("UPDATE conversations SET preview = ?, updated_at = ? WHERE id = ?", (payload.content[:60], now(), conversation_id))
                assistant_seq = message_count + 1
                for attachment in email_attachments:
                    conn.execute("INSERT OR REPLACE INTO conversation_attachments (id, conversation_id, message_seq, attachment_id, filename, size, content_type, provider_message_id, account_id, created_at) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)", (f"CA-{uuid.uuid4().hex[:12].upper()}", conversation_id, assistant_seq, str(attachment.get('id') or attachment.get('name')), str(attachment.get('name') or '附件'), int(attachment.get('size') or 0), str(attachment.get('contentType') or 'application/octet-stream'), remote_id, account_id, now()))
                audit(conn, user['id'], 'email.read', 'conversation', conversation_id, {'task_id': task_id, 'count': len(provider_items), 'detail_requested': detail_requested, 'index': requested_index})
                conn.commit()
            return {'conversation_id': conversation_id, 'task_id': task_id, 'trace_id': trace_id, 'response': response_text, 'attachments': email_attachments, 'events': _conversation_task_events(conversation_id, task_id)}
        except RuntimeError as exc:
            with closing(connect()) as conn:
                _insert_agent_event(conn, conversation_id, task_id, trace_id, event_seq + 2, 'tool.email_list', '读取企业邮箱', 'failed', '--', str(exc), {'error': str(exc)})
                audit(conn, user['id'], 'email.read_failed', 'conversation', conversation_id, {'task_id': task_id, 'error': str(exc)})
                conn.commit()
            raise HTTPException(status_code=502, detail=f'邮件读取失败：{exc}') from exc

    if email_kind == 'send':
        recipients, subject, body = _parse_email_send_request(payload.content)
        if not recipients or not subject or not body:
            response_text = '我可以帮你发送邮件。请按以下格式补充：给 name@example.com 发邮件，主题：会议安排，正文：请查看附件。'
        else:
            pending = {'to': recipients, 'subject': subject, 'body': body, 'account_id': _active_account_id(user)}
            response_text = '邮件内容已准备好：\n收件人：' + '、'.join(recipients) + '\n主题：' + subject + '\n正文：' + body + '\n\n确认发送吗？请回复“确认发送”。'
        with closing(connect()) as conn:
            status = 'pending' if recipients and subject and body else 'success'
            event_type = 'email.send_pending' if status == 'pending' else 'email.send_validate'
            summary = '等待用户确认后发送' if status == 'pending' else '缺少收件人、主题或正文'
            _insert_agent_event(conn, conversation_id, task_id, trace_id, event_seq + 2, event_type, '准备发送企业邮件', status, '已完成', summary, pending if status == 'pending' else {'missing': True})
            conn.execute("INSERT INTO messages VALUES (?, ?, 'assistant', ?, ?, ?)", (f'MSG-{uuid.uuid4().hex[:10].upper()}', conversation_id, response_text, message_count + 1, now()))
            conn.execute("UPDATE conversations SET preview = ?, updated_at = ? WHERE id = ?", (payload.content[:60], now(), conversation_id))
            audit(conn, user['id'], 'email.send_prepare', 'conversation', conversation_id, {'task_id': task_id, 'recipient_count': len(recipients), 'awaiting_confirmation': status == 'pending'})
            conn.commit()
        return {'conversation_id': conversation_id, 'task_id': task_id, 'trace_id': trace_id, 'response': response_text, 'events': _conversation_task_events(conversation_id, task_id)}
    if document_request:
        try:
            request_text = chr(10).join(str(item["content"]) for item in history if item["role"] == "user")
            leave_details = _extract_leave_details(request_text)
            is_leave_request = any(keyword in request_text for keyword in ("请假", "病假", "年假", "事假", "调休"))
            leave_request_id = None
            if is_leave_request:
                missing = [label for key, label in (("start", "开始日期"), ("end", "结束日期"), ("reason", "申请事由")) if leave_details[key] == "未填写"]
                with closing(connect()) as conn:
                    _insert_agent_event(conn, conversation_id, task_id, trace_id, event_seq + 2, "leave.parse", "解析请假字段", "success", "已完成", f"已识别：{leave_details['leave_type']}，{leave_details['start']} 至 {leave_details['end']}，{leave_details['days']} 天", leave_details)
                    if missing:
                        response_text = "已识别请假请求，但还缺少：" + "、".join(missing) + "。请补充后再提交，系统不会创建审批申请或文件。"
                        _insert_agent_event(conn, conversation_id, task_id, trace_id, event_seq + 3, "leave.validate", "校验请假信息", "failed", "已完成", "缺少必填信息：" + "、".join(missing), {"missing": missing})
                        conn.execute("INSERT INTO messages VALUES (?, ?, 'assistant', ?, ?, ?)", (f"MSG-{uuid.uuid4().hex[:10].upper()}", conversation_id, response_text, message_count + 1, now()))
                        conn.execute("UPDATE conversations SET preview = ?, updated_at = ? WHERE id = ?", (payload.content[:60], now(), conversation_id))
                        audit(conn, user["id"], "leave.validation_failed", "conversation", conversation_id, {"missing": missing})
                        conn.commit()
                        events = rows(conn, "SELECT id, task_id, trace_id, seq, label, event_type AS type, status, duration, safe_summary AS summary, data, created_at FROM agent_events WHERE conversation_id = ? AND task_id = ? ORDER BY seq", (conversation_id, task_id))
                        return {"conversation_id": conversation_id, "task_id": task_id, "trace_id": trace_id, "response": response_text, "events": events}
                    leave_request_id, _ = _ensure_leave_request(conn, user, leave_details)
                    _insert_agent_event(conn, conversation_id, task_id, trace_id, event_seq + 3, "tool.leave_request_create", "创建待审批申请", "success", "已完成", f"已创建申请单 {leave_request_id}，已分配直属上级审批", {"leave_request_id": leave_request_id, **leave_details})
                    conn.commit()
            file_id, file_path = _create_docx_file(conversation_id, user, request_text, leave_request_id=leave_request_id)
            filename = Path(file_path).name
            with closing(connect()) as conn:
                _insert_agent_event(conn, conversation_id, task_id, trace_id, event_seq + 4, "tool.document_generate", "生成 Office 文档", "success", "已完成", "已创建 DOCX 文件并关联审批申请", {"file_id": file_id, "filename": filename, "leave_request_id": leave_request_id, "download_url": f"/api/files/{file_id}/download"})
                response_text = f"已生成 DOCX 文档：{filename}。关联申请单：{leave_request_id or '无'}。文件已保存，可在“文件中心”中查看并下载。"
                conn.execute("INSERT INTO generated_files (id, name, file_type, status, conversation_id, created_at, template, file_path, owner_id) VALUES (?, ?, 'docx', 'success', ?, ?, ?, ?, ?)", (file_id, filename, conversation_id, now(), "企业服务申请", file_path, user["id"]))
                conn.execute("UPDATE conversations SET preview = ?, updated_at = ? WHERE id = ?", (payload.content[:60], now(), conversation_id))
                conn.execute("INSERT INTO messages VALUES (?, ?, 'assistant', ?, ?, ?)", (f"MSG-{uuid.uuid4().hex[:10].upper()}", conversation_id, response_text, message_count + 1, now()))
                audit(conn, user["id"], "file.generate", "generated_file", file_id, {"conversation_id": conversation_id, "leave_request_id": leave_request_id, "source": "agent", "request": request_text})
                conn.commit()
                events = rows(conn, "SELECT id, task_id, trace_id, seq, label, event_type AS type, status, duration, safe_summary AS summary, data, created_at FROM agent_events WHERE conversation_id = ? AND task_id = ? ORDER BY seq", (conversation_id, task_id))
            return {"conversation_id": conversation_id, "task_id": task_id, "trace_id": trace_id, "response": response_text, "events": events, "generated_file_id": file_id, "download_url": f"/api/files/{file_id}/download"}
        except Exception as exc:
            with closing(connect()) as conn:
                _insert_agent_event(conn, conversation_id, task_id, trace_id, event_seq + 4, "tool.document_generate", "生成 Office 文档", "failed", "--", str(exc), {"error": str(exc)})
                audit(conn, user["id"], "file.generate_failed", "conversation", conversation_id, {"task_id": task_id, "error": str(exc)})
                conn.commit()
            raise HTTPException(status_code=500, detail=f"文档生成失败：{exc}") from exc
    try:
        response_text, duration_ms, endpoint, usage = call_configured_model(config, [{"role": "system", "content": "你是企业内部智能服务助手。必须结合当前会话历史理解“继续、它、这个”等指代，不能把每条消息当成孤立问题。优先准确回答，不要编造制度；涉及提交、发送、删除等写操作时，只能说明需要确认。"}, *history])
    except ValueError as exc:
        with closing(connect()) as conn:
            _insert_agent_event(conn, conversation_id, task_id, trace_id, event_seq + 2, "model_requested", "调用已配置模型", "failed", "--", str(exc), {"model": config.get("model", "unknown"), "error": str(exc)})
            audit(conn, user["id"], "agent.model_failed", "conversation", conversation_id, {"task_id": task_id, "trace_id": trace_id})
            conn.commit()
        raise HTTPException(502, str(exc)) from exc
    with closing(connect()) as conn:
        _insert_agent_event(conn, conversation_id, task_id, trace_id, event_seq + 2, "model_requested", "调用已配置模型", "success", f"{duration_ms}ms", "模型调用成功", {"model": config["model"], "endpoint": endpoint, "response_saved": True, "usage": usage})
        _insert_agent_event(conn, conversation_id, task_id, trace_id, event_seq + 3, "response_generated", "生成结果", "success", f"{duration_ms}ms", "已生成模型回答", {"model": config["model"], "citation_count": 0, "usage": usage})
        conn.execute("INSERT INTO token_usage (id, user_id, conversation_id, task_id, provider, model, input_tokens, output_tokens, total_tokens, created_at) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?)", (f"USE-{uuid.uuid4().hex[:12].upper()}", user["id"], conversation_id, task_id, config.get("provider", ""), config["model"], usage["input_tokens"], usage["output_tokens"], usage["total_tokens"], now()))
        conn.execute("INSERT INTO messages VALUES (?, ?, 'assistant', ?, ?, ?)", (f"MSG-{uuid.uuid4().hex[:10].upper()}", conversation_id, response_text, message_count + 1, now()))
        conn.execute("UPDATE conversations SET preview = ?, updated_at = ? WHERE id = ?", (payload.content[:60], now(), conversation_id))
        audit(conn, user["id"], "conversation.message", "conversation", conversation_id, {"task_id": task_id, "trace_id": trace_id, "model": config["model"]})
        conn.commit()
        events = rows(conn, "SELECT id, seq, label, event_type AS type, status, duration, safe_summary AS summary, data FROM agent_events WHERE conversation_id = ? AND task_id = ? ORDER BY seq", (conversation_id, task_id))
    return {"conversation_id": conversation_id, "task_id": task_id, "trace_id": trace_id, "response": response_text, "events": events}

@app.get("/api/conversations/{conversation_id}/trace")
def conversation_trace(conversation_id: str, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    require_role(user, "manager", "admin")
    with closing(connect()) as conn:
        conversation = row(conn, "SELECT id, title, preview, status, updated_at AS updated FROM conversations WHERE id = ?", (conversation_id,))
        if not conversation:
            raise HTTPException(404, "conversation not found")
        messages = rows(conn, "SELECT id, role, content, seq, created_at FROM messages WHERE conversation_id = ? ORDER BY seq", (conversation_id,))
        events = rows(conn, "SELECT id, task_id, trace_id, seq, label, event_type AS type, status, duration, safe_summary AS summary, data, created_at FROM agent_events WHERE conversation_id = ? ORDER BY created_at, seq", (conversation_id,))
        audit_items = rows(conn, "SELECT id, action, resource_type, resource_id, detail, created_at FROM audit_events WHERE resource_id = ? OR resource_id IN (SELECT id FROM messages WHERE conversation_id = ?) ORDER BY created_at", (conversation_id, conversation_id))
    return {"conversation": conversation, "messages": messages, "events": events, "audit": audit_items}


def _send_agentmail_message(payload: AgentMailSendRequest, user: dict[str, Any], source: str, attachment_paths: list[Path] | None = None, attachment_meta: list[dict[str, Any]] | None = None) -> dict[str, Any]:
    if not payload.confirmed:
        raise HTTPException(status_code=409, detail='发送邮件前需要用户确认')
    if not payload.subject.strip() and not payload.body.strip():
        raise HTTPException(status_code=422, detail='邮件主题或正文不能为空')
    arguments = ['message', '+send']
    for address in payload.to:
        if address.strip():
            arguments.extend(['--to', address.strip()])
    for address in payload.cc:
        if address.strip():
            arguments.extend(['--cc', address.strip()])
    for address in payload.bcc:
        if address.strip():
            arguments.extend(['--bcc', address.strip()])
    arguments.extend(['--subject', payload.subject, '--body', payload.body, '--confirmed'])
    for attachment_path in attachment_paths or []:
        arguments.extend(['--attachment', str(attachment_path)])
    workspace = _account_workspace(user, payload.account_id)
    data = _run_agentmail(arguments, workspace=workspace)
    provider_message = data.get('message') if isinstance(data.get('message'), dict) else data
    if not isinstance(provider_message, dict):
        provider_message = {'subject': payload.subject, 'to': [{'email': item} for item in payload.to]}
    if attachment_meta and not provider_message.get('attachments'):
        provider_message['attachments'] = attachment_meta
    with closing(connect()) as conn:
        saved = _upsert_email(conn, user['id'], provider_message, 'sent', source, payload.body, _active_account_id(user, payload.account_id))
        audit(conn, user['id'], 'email.send', 'email_message', saved.get('id'), {'provider_id': saved.get('provider_id'), 'source': source, 'to': payload.to})
        conn.commit()
    return _email_response(saved)


@app.get('/api/email/accounts')
def list_email_accounts(user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    with closing(connect()) as conn:
        accounts = rows(conn, "SELECT * FROM email_accounts WHERE owner_id = ? AND status = 'active' ORDER BY is_active DESC, updated_at DESC", (user['id'],))
    if not accounts:
        try:
            profile = _run_agentmail(['+me'])
            accounts = [_save_agent_account(user['id'], 'default', profile, '默认 AgentMail 邮箱')]
        except RuntimeError:
            accounts = []
    else:
        accounts = [_account_public(item) for item in accounts]
    return {'accounts': accounts}


@app.post('/api/email/accounts/auth')
def start_email_account_auth(payload: EmailAuthStartRequest, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    workspace = (payload.workspace or f"platform-{user['id']}-{uuid.uuid4().hex[:8]}").strip()
    if not re.fullmatch(r'[A-Za-z0-9_-]{3,80}', workspace):
        raise HTTPException(422, '邮箱工作区标识格式无效')
    job_id = f"EAJ-{uuid.uuid4().hex[:12].upper()}"
    with EMAIL_AUTH_LOCK:
        EMAIL_AUTH_JOBS[job_id] = {'status': 'waiting', 'workspace': workspace, 'owner_id': user['id'], 'output': '', 'authorization_url': ''}
    threading.Thread(target=_run_email_auth, args=(job_id, user['id'], workspace, payload.label.strip()), daemon=True).start()
    return {'job_id': job_id, 'status': 'waiting', 'workspace': workspace, 'message': '授权流程已启动，请在新打开的浏览器窗口中使用微信扫码完成 AgentMail 授权。'}


@app.get('/api/email/accounts/auth/{job_id}')
def email_account_auth_status(job_id: str, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    with EMAIL_AUTH_LOCK:
        job = dict(EMAIL_AUTH_JOBS.get(job_id) or {})
    if not job or job.get('owner_id') != user['id']:
        raise HTTPException(404, '授权任务不存在或已过期')
    return {key: job.get(key) for key in ('status', 'workspace', 'authorization_url', 'output', 'message', 'account')}


@app.post('/api/email/accounts/{account_id}/activate')
def activate_email_account(account_id: str, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    stamp = now()
    with closing(connect()) as conn:
        account = row(conn, "SELECT * FROM email_accounts WHERE id = ? AND owner_id = ? AND status = 'active'", (account_id, user['id']))
        if not account:
            raise HTTPException(404, '邮箱账号不存在或无权使用')
        conn.execute("UPDATE email_accounts SET is_active = CASE WHEN id = ? THEN 1 ELSE 0 END, updated_at = ? WHERE owner_id = ? AND status = 'active'", (account_id, stamp, user['id']))
        conn.commit()
        updated = row(conn, "SELECT * FROM email_accounts WHERE id = ?", (account_id,))
    return _account_public(updated or account)

@app.get('/api/email/attachments/{message_id}/{attachment_id}')
def download_email_attachment(message_id: str, attachment_id: str, account_id: str | None = None, conversation_id: str | None = None, user: dict[str, Any] = Depends(get_user)) -> FileResponse:
    account_ref = _active_account_id(user, account_id)
    output_dir = EMAIL_ATTACHMENT_DIR / 'downloads'
    output_dir.mkdir(parents=True, exist_ok=True)
    try:
        data = _run_agentmail(['attachment', '+download', '--msg', message_id, '--att', attachment_id, '--output', str(EMAIL_ATTACHMENT_DIR / 'downloads')], workspace=_account_workspace(user, account_ref))
    except RuntimeError as exc:
        raise HTTPException(status_code=502, detail=f'附件下载失败：{exc}') from exc
    saved_to = data.get('saved_to') or data.get('path') or data.get('file_path') if isinstance(data, dict) else None
    candidate = Path(str(saved_to)) if saved_to else None
    if candidate and not candidate.is_absolute():
        candidate = ROOT / candidate
    if not candidate or not candidate.exists():
        matches = sorted(output_dir.glob('*'), key=lambda item: item.stat().st_mtime, reverse=True)
        candidate = matches[0] if matches else None
    if not candidate or not candidate.exists() or output_dir.resolve() not in candidate.resolve().parents:
        raise HTTPException(status_code=404, detail='附件文件未找到')
    with closing(connect()) as conn:
        audit(conn, user['id'], 'email.attachment_download', 'email_attachment', attachment_id, {'message_id': message_id, 'conversation_id': conversation_id, 'account_id': account_ref})
        conn.commit()
    return FileResponse(candidate, filename=candidate.name)
@app.get('/api/email/messages')
def list_email_messages(dir: Literal['inbox', 'sent', 'trash'] = 'inbox', account_id: str | None = None, user: dict[str, Any] = Depends(get_user)) -> list[dict[str, Any]]:
    provider_items: list[dict[str, Any]] = []
    try:
        data = _run_agentmail(['message', '+list', '--dir', dir, '--limit', '50'], workspace=_account_workspace(user, account_id))
        provider_items = data.get('data') if isinstance(data.get('data'), list) else []
        with closing(connect()) as conn:
            for item in provider_items:
                if isinstance(item, dict):
                    _upsert_email(conn, user['id'], item, dir, 'sync', account_id=_active_account_id(user, account_id))
            conn.commit()
    except RuntimeError:
        pass
    with closing(connect()) as conn:
        cached = rows(conn, "SELECT id, provider_id, owner_id, account_id, folder, direction, sender_email, sender_name, recipients, cc, subject, preview, '' AS body, is_read, has_attachments, attachments, source, status, provider_payload, created_at, updated_at FROM email_messages WHERE owner_id = ? AND account_id = ? AND folder = ? ORDER BY created_at DESC", (user['id'], _active_account_id(user, account_id), dir))
    if not cached and not provider_items:
        raise HTTPException(status_code=502, detail='AgentMail 邮件服务暂不可用')
    return [_email_response(item) for item in cached]


@app.post('/api/email/sync')
def sync_email_messages(dir: Literal['inbox', 'sent', 'trash'] = 'inbox', account_id: str | None = None, user: dict[str, Any] = Depends(get_user)) -> list[dict[str, Any]]:
    data = _run_agentmail(['message', '+list', '--dir', dir, '--limit', '50'], workspace=_account_workspace(user, account_id))
    provider_items = data.get('data') if isinstance(data.get('data'), list) else []
    with closing(connect()) as conn:
        for item in provider_items:
            if isinstance(item, dict):
                _upsert_email(conn, user['id'], item, dir, 'sync', account_id=_active_account_id(user, account_id))
        conn.commit()
        cached = rows(conn, "SELECT id, provider_id, owner_id, account_id, folder, direction, sender_email, sender_name, recipients, cc, subject, preview, '' AS body, is_read, has_attachments, attachments, source, status, provider_payload, created_at, updated_at FROM email_messages WHERE owner_id = ? AND account_id = ? AND folder = ? ORDER BY created_at DESC", (user['id'], _active_account_id(user, account_id), dir))
    return [_email_response(item) for item in cached]


@app.get('/api/email/messages/{provider_id}')
def read_email_message(provider_id: str, account_id: str | None = None, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    data = _run_agentmail(['message', '+read', '--id', provider_id], workspace=_account_workspace(user, account_id))
    provider_message = data.get('message') if isinstance(data.get('message'), dict) else data
    if not isinstance(provider_message, dict):
        raise HTTPException(status_code=502, detail='AgentMail 返回的邮件格式无法识别')
    with closing(connect()) as conn:
        account_ref = _active_account_id(user, account_id)
        stored_provider_id = _mail_provider_key(provider_id, account_ref)
        previous = row(conn, 'SELECT folder FROM email_messages WHERE owner_id = ? AND provider_id = ?', (user['id'], stored_provider_id))
        saved = _upsert_email(conn, user['id'], provider_message, previous.get('folder', 'inbox') if previous else 'inbox', 'read', account_id=account_ref)
        conn.execute('UPDATE email_messages SET is_read = 1 WHERE provider_id = ?', (stored_provider_id,))
        conn.commit()
    saved['is_read'] = 1
    return _email_response(saved)


@app.post('/api/email/send', status_code=201)
def send_email(payload: AgentMailSendRequest, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    return _send_agentmail_message(payload, user, 'visual')


@app.post('/api/email/send-with-attachments', status_code=201)
async def send_email_with_attachments(
    to: str = Form(...),
    subject: str = Form(''),
    body: str = Form(''),
    confirmed: bool = Form(False),
    account_id: str | None = Form(None),
    attachments: list[UploadFile] = File(default=[]),
    user: dict[str, Any] = Depends(get_user),
) -> dict[str, Any]:
    try:
        recipients = json.loads(to)
        if not isinstance(recipients, list):
            raise ValueError
        recipients = [str(item).strip() for item in recipients if str(item).strip()]
    except (json.JSONDecodeError, ValueError):
        recipients = [item.strip() for item in re.split(r'[,;\s]+', to) if item.strip()]
    if not recipients:
        raise HTTPException(status_code=422, detail='请至少填写一个收件人')
    if len(attachments) > 10:
        raise HTTPException(status_code=422, detail='单封邮件最多添加 10 个附件')
    paths: list[Path] = []
    metadata: list[dict[str, Any]] = []
    for upload in attachments:
        filename = Path(upload.filename or '附件').name
        if not filename or filename in {'.', '..'}:
            raise HTTPException(status_code=422, detail='附件文件名无效')
        content = await upload.read()
        if len(content) > 25 * 1024 * 1024:
            raise HTTPException(status_code=413, detail=f'附件 {filename} 超过 25 MB 限制')
        path = EMAIL_ATTACHMENT_DIR / f'{uuid.uuid4().hex}_{filename}'
        path.write_bytes(content)
        paths.append(path)
        metadata.append({'filename': filename, 'size': len(content), 'content_type': upload.content_type or 'application/octet-stream', 'source': 'uploaded'})
    payload = AgentMailSendRequest(to=recipients, subject=subject, body=body, confirmed=confirmed, account_id=account_id)
    try:
        return _send_agentmail_message(payload, user, 'visual', paths, metadata)
    except RuntimeError as exc:
        raise HTTPException(status_code=502, detail=str(exc)) from exc

@app.post('/api/email/agent/send', status_code=201)
def agent_send_email(payload: AgentMailSendRequest, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    return _send_agentmail_message(payload, user, 'agent')

def _file_record(file_id: str, user: dict[str, Any]) -> dict[str, Any]:
    with closing(connect()) as conn:
        if user["role"] in {"manager", "admin"}:
            item = row(conn, "SELECT * FROM generated_files WHERE id = ?", (file_id,))
        else:
            item = row(conn, "SELECT * FROM generated_files WHERE id = ? AND (owner_id = ? OR (owner_id = '' AND conversation_id IN (SELECT id FROM conversations WHERE user_id = ?)))", (file_id, user["id"], user["id"]))
    if not item:
        raise HTTPException(404, "file not found")
    if not item.get("file_path") or not Path(item["file_path"]).exists():
        raise HTTPException(404, "file content not available")
    return item


def _file_media_type(name: str) -> str:
    suffix = Path(name).suffix.lower()
    return {
        ".pdf": "application/pdf",
        ".doc": "application/msword",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xls": "application/vnd.ms-excel",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".ppt": "application/vnd.ms-powerpoint",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".json": "application/json",
        ".csv": "text/csv",
    }.get(suffix, "application/octet-stream")


def _extract_xlsx_text(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(f"[工作表] {sheet.title}")
        for values in sheet.iter_rows(values_only=True):
            cells = ["" if value is None else str(value) for value in values]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx_text(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(path)
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"[幻灯片 {index}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n".join(parts)

def _extract_docx_text(path: Path) -> str:
    from docx import Document
    document = Document(path)
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for table_row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in table_row.cells))
    return "\n".join(parts)


@app.get("/api/files/{file_id}/download")
def download_file(file_id: str, user: dict[str, Any] = Depends(get_user)) -> FileResponse:
    item = _file_record(file_id, user)
    return FileResponse(item["file_path"], filename=item["name"], media_type=_file_media_type(item["name"]))


@app.get("/api/files/{file_id}/preview")
def preview_file(file_id: str, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    item = _file_record(file_id, user)
    path = Path(item["file_path"])
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return {"kind": "pdf", "name": item["name"], "download_url": f"/api/files/{file_id}/download"}
    if suffix == ".doc":
        raise HTTPException(415, "暂不支持旧版 .doc 在线解析，请下载后查看")
    if suffix == ".docx":
        content = _extract_docx_text(path)
    elif suffix in {".txt", ".md", ".csv", ".json", ".xml", ".log"}:
        content = path.read_text(encoding="utf-8", errors="replace")
    else:
        raise HTTPException(415, "该文件类型暂不支持在线预览，请下载查看")
    return {"kind": "text", "name": item["name"], "content": content}


@app.get("/api/files")
def list_files(user: dict[str, Any] = Depends(get_user)) -> list[dict[str, Any]]:
    with closing(connect()) as conn:
        if user["role"] in {"manager", "admin"}:
            return rows(conn, "SELECT id, name, file_type AS type, status, conversation_id AS conversationId, created_at AS createdAt, template, file_path AS filePath, owner_id FROM generated_files ORDER BY created_at DESC")
        return rows(conn, "SELECT id, name, file_type AS type, status, conversation_id AS conversationId, created_at AS createdAt, template, file_path AS filePath, owner_id FROM generated_files WHERE owner_id = ? OR (owner_id = '' AND conversation_id IN (SELECT id FROM conversations WHERE user_id = ?)) ORDER BY created_at DESC", (user["id"], user["id"]))

@app.post("/api/files/upload", status_code=201)
async def upload_file(file: UploadFile = File(...), user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    if not file.filename:
        raise HTTPException(422, "文件名不能为空")
    suffix = Path(file.filename).suffix.lower()
    allowed = {".pdf", ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx", ".md", ".txt", ".csv", ".json", ".xml", ".log"}
    if suffix not in allowed:
        raise HTTPException(415, "暂不支持该文件类型")
    content = await file.read()
    if not content:
        raise HTTPException(422, "不能上传空文件")
    file_id = f"FILE-{uuid.uuid4().hex[:10].upper()}"
    output_dir = DATA_DIR / "files"
    output_dir.mkdir(parents=True, exist_ok=True)
    path = output_dir / f"{file_id}{suffix}"
    path.write_bytes(content)
    stamp = now()
    with closing(connect()) as conn:
        conn.execute("INSERT INTO generated_files (id, name, file_type, status, conversation_id, created_at, template, file_path, owner_id) VALUES (?, ?, ?, 'success', 'FILE-CENTER', ?, '文件上传', ?, ?)", (file_id, file.filename, suffix.lstrip(".") or "file", stamp, str(path), user["id"]))
        audit(conn, user["id"], "file.upload", "file", file_id, {"name": file.filename, "size": len(content), "type": suffix})
        conn.commit()
        saved = row(conn, "SELECT id, name, file_type AS type, status, conversation_id AS conversationId, created_at AS createdAt, template, file_path AS filePath, owner_id FROM generated_files WHERE id = ?", (file_id,))
    return saved or {}

@app.post("/api/files/generate", status_code=201)
def generate_file(conversation_id: str, template: str = "员工请假条", user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    file_id = f"FILE-{uuid.uuid4().hex[:8].upper()}"
    name = f"{user['name']}-请假条.docx" if "请假" in template else "企业服务生成文件.docx"
    with closing(connect()) as conn:
        conn.execute("INSERT INTO generated_files (id, name, file_type, status, conversation_id, created_at, template, file_path, owner_id) VALUES (?, ?, 'docx', 'processing', ?, '刚刚', ?, '', ?)", (file_id, name, conversation_id, template, user["id"]))
        audit(conn, user["id"], "file.generate", "generated_file", file_id, {"conversation_id": conversation_id, "template": template})
        conn.commit()
    return {"id": file_id, "name": name, "status": "processing", "conversationId": conversation_id, "template": template}


def model_endpoints(api_url: str) -> list[str]:
    base = api_url.strip().rstrip('/')
    parsed = urllib.parse.urlparse(base)
    if parsed.scheme not in {"http", "https"} or not parsed.netloc:
        raise HTTPException(400, "API URL must be an absolute http(s) URL")
    if parsed.path.rstrip('/').endswith('/models'):
        return [base]
    candidates = [f"{base}/models", f"{base}/v1/models"]
    if parsed.path.rstrip('/').endswith('/v1'):
        candidates = [f"{base}/models", f"{base[:-3].rstrip('/')}/models"]
    return list(dict.fromkeys(candidates))


def extract_model_names(payload: Any) -> list[str]:
    items = payload.get("data") if isinstance(payload, dict) else payload
    if items is None and isinstance(payload, dict):
        items = payload.get("models")
    if not isinstance(items, list):
        return []
    names: list[str] = []
    for item in items:
        if isinstance(item, str):
            value = item
        elif isinstance(item, dict):
            value = item.get("id") or item.get("name") or item.get("model")
        else:
            value = None
        if isinstance(value, str) and value.strip() and value.strip() not in names:
            names.append(value.strip())
    return names


@app.post("/api/model-configs/discover-models")
def discover_models(payload: ModelDiscoveryRequest, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    require_role(user, "admin")
    try:
        model_names, endpoint = fetch_models_from_provider(
            payload.api_url,
            payload.api_key,
            api_format=payload.api_format,
            is_full_url=payload.is_full_url,
            models_url_override=payload.models_url_override,
        )
    except ValueError as exc:
        raise HTTPException(400, str(exc)) from exc

    with closing(connect()) as conn:
        audit(conn, user["id"], "model_config.discover", "model_config", None, {"provider": payload.provider, "endpoint": endpoint, "model_count": len(model_names)})
        conn.commit()
    return {"models": model_names, "endpoint": endpoint}

@app.get("/api/model-configs")
def list_model_configs(user: dict[str, Any] = Depends(get_user)) -> list[dict[str, Any]]:
    require_role(user, "admin")
    with closing(connect()) as conn:
        return rows(conn, "SELECT id, scope_type, scope_id, provider, model, api_url, api_key_masked, api_format, enabled, updated_at FROM model_configs ORDER BY scope_type, scope_id")


@app.put("/api/model-configs")
def save_model_config(payload: ModelConfigRequest, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    require_role(user, "admin")
    config_id = f"MC-{payload.scope_type}-{payload.scope_id}"
    with closing(connect()) as conn:
        existing = row(conn, "SELECT api_key_ciphertext FROM model_configs WHERE scope_type = ? AND scope_id = ?", (payload.scope_type, payload.scope_id))
        ciphertext = encrypt_secret(payload.api_key.strip()) if payload.api_key.strip() else (existing or {}).get("api_key_ciphertext", "")
        masked = "sk-••••••••" if ciphertext else "未配置"
        conn.execute("INSERT INTO model_configs (id, scope_type, scope_id, provider, model, api_url, api_key_masked, api_format, api_key_ciphertext, enabled, updated_at) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?) ON CONFLICT(scope_type, scope_id) DO UPDATE SET provider=excluded.provider, model=excluded.model, api_url=excluded.api_url, api_key_masked=excluded.api_key_masked, api_format=excluded.api_format, api_key_ciphertext=excluded.api_key_ciphertext, enabled=excluded.enabled, updated_at=excluded.updated_at", (config_id, payload.scope_type, payload.scope_id, payload.provider, payload.model, payload.api_url, masked, payload.api_format, ciphertext, int(payload.enabled), now()))
        audit(conn, user["id"], "model_config.update", "model_config", config_id, {"scope_type": payload.scope_type, "scope_id": payload.scope_id, "api_format": payload.api_format})
        conn.commit()
        return row(conn, "SELECT id, scope_type, scope_id, provider, model, api_url, api_key_masked, api_format, enabled, updated_at FROM model_configs WHERE id = ?", (config_id,)) or {}

@app.get("/api/usage")
def usage_summary(day: str | None = None, provider: str | None = None, model: str | None = None, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    if day:
        try:
            datetime.strptime(day, "%Y-%m-%d")
        except ValueError as exc:
            raise HTTPException(422, "日期格式必须为 YYYY-MM-DD") from exc
    with closing(connect()) as conn:
        if user["role"] == "admin":
            scope = rows(conn, "SELECT id, name, email, department, role FROM users WHERE status = 'active' ORDER BY department, name")
        elif user["role"] == "manager":
            scope = rows(conn, """SELECT id, name, email, department, role FROM users
                WHERE status = 'active' AND (id = ? OR supervisor_id = ? OR department IN
                (SELECT name FROM departments WHERE manager_id = ?))
                ORDER BY department, name""", (user["id"], user["id"], user["id"]))
        else:
            scope = [user]
        ids = [item["id"] for item in scope]
        if not ids:
            return {"scope": user["role"], "selectedDay": day, "filters": {"providers": [], "models": []}, "summary": {"inputTokens": 0, "outputTokens": 0, "totalTokens": 0, "requests": 0}, "trend": [], "members": []}
        placeholders = ",".join("?" for _ in ids)
        base_params: list[Any] = list(ids)
        filters = [f"user_id IN ({placeholders})"]
        if day:
            filters.append("substr(created_at, 1, 10) = ?")
            base_params.append(day)
        if provider:
            filters.append("provider = ?")
            base_params.append(provider)
        if model:
            filters.append("model = ?")
            base_params.append(model)
        where = " AND ".join(filters)
        available_rows = rows(conn, f"SELECT DISTINCT provider, model FROM token_usage WHERE user_id IN ({placeholders}) ORDER BY provider, model", tuple(ids))
        aggregates = rows(conn, f"""SELECT user_id AS userId,
            COALESCE(SUM(input_tokens), 0) AS inputTokens,
            COALESCE(SUM(output_tokens), 0) AS outputTokens,
            COALESCE(SUM(total_tokens), 0) AS totalTokens,
            COUNT(*) AS requests
            FROM token_usage WHERE {where} GROUP BY user_id""", tuple(base_params))
        trend_rows = rows(conn, f"""SELECT substr(created_at, 12, 2) AS hour,
            COALESCE(SUM(input_tokens), 0) AS inputTokens,
            COALESCE(SUM(output_tokens), 0) AS outputTokens,
            COALESCE(SUM(total_tokens), 0) AS totalTokens,
            COUNT(*) AS requests
            FROM token_usage WHERE {where} GROUP BY substr(created_at, 12, 2) ORDER BY hour""", tuple(base_params))
        by_user = {item["userId"]: item for item in aggregates}
        members = []
        for item in scope:
            usage = by_user.get(item["id"], {"inputTokens": 0, "outputTokens": 0, "totalTokens": 0, "requests": 0})
            members.append({**item, **usage, "userId": item["id"]})
        summary = {key: sum(int(item[key]) for item in members) for key in ("inputTokens", "outputTokens", "totalTokens")}
        summary["requests"] = sum(int(item["requests"]) for item in members)
        by_hour = {int(item["hour"]): item for item in trend_rows if str(item.get("hour", "")).isdigit()}
        trend = [{"hour": hour, "label": f"{hour:02d}:00", "inputTokens": int(by_hour.get(hour, {}).get("inputTokens", 0)), "outputTokens": int(by_hour.get(hour, {}).get("outputTokens", 0)), "totalTokens": int(by_hour.get(hour, {}).get("totalTokens", 0)), "requests": int(by_hour.get(hour, {}).get("requests", 0))} for hour in range(24)]
    return {
        "scope": user["role"],
        "selectedDay": day,
        "filters": {"providers": sorted({str(item.get("provider") or "") for item in available_rows if item.get("provider")}), "models": sorted({str(item.get("model") or "") for item in available_rows if item.get("model")})},
        "summary": summary,
        "trend": trend,
        "members": members,
    }
@app.get("/api/monitoring")
def monitoring(user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    require_role(user, "manager", "admin")
    with closing(connect()) as conn:
        active = row(conn, "SELECT COUNT(*) AS count FROM ingestion_tasks WHERE status = 'processing'")["count"]
        failed = row(conn, "SELECT COUNT(*) AS count FROM ingestion_tasks WHERE status = 'failed'")["count"]
        total_tasks = row(conn, "SELECT COUNT(*) AS count FROM ingestion_tasks")["count"]
        successful = row(conn, "SELECT COUNT(*) AS count FROM ingestion_tasks WHERE status = 'success'")["count"]
    rate = f"{successful / total_tasks * 100:.1f}%" if total_tasks else "--"
    return {"total_tasks": total_tasks, "success_rate": rate, "active_ingestion": active, "failed_ingestion": failed, "agent_success_rate": "--"}


@app.get("/api/admin/users")
def admin_users(user: dict[str, Any] = Depends(get_user)) -> list[dict[str, Any]]:
    require_role(user, "admin")
    with closing(connect()) as conn:
        return rows(conn, "SELECT id, name, email, department, role, status FROM users ORDER BY department, name")


@app.post("/api/admin/users", status_code=201)
def create_admin_user(payload: CreateUserRequest, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    require_role(user, "admin")
    email = payload.email.strip().lower()
    if "@" not in email:
        raise HTTPException(422, "请输入有效邮箱")
    user_id = f"usr-{uuid.uuid4().hex[:10]}"
    with closing(connect()) as conn:
        if row(conn, "SELECT id FROM users WHERE lower(email) = lower(?)", (email,)):
            raise HTTPException(409, "该邮箱已存在")
        if not row(conn, "SELECT id FROM departments WHERE name = ?", (payload.department,)):
            raise HTTPException(422, "部门不存在，请先创建或选择有效部门")
        conn.execute(
            "INSERT INTO users (id, name, email, department, role, status, password_hash) VALUES (?, ?, ?, ?, ?, 'active', ?)",
            (user_id, payload.name.strip(), email, payload.department, payload.role, hash_password(payload.password)),
        )
        audit(conn, user["id"], "user.create", "user", user_id, {"email": email, "department": payload.department, "role": payload.role})
        conn.commit()
        created = row(conn, "SELECT id, name, email, department, role, status FROM users WHERE id = ?", (user_id,))
    return created or {}

@app.get("/api/admin/password-reset-requests")
def admin_password_reset_requests(user: dict[str, Any] = Depends(get_user)) -> list[dict[str, Any]]:
    require_role(user, "admin")
    with closing(connect()) as conn:
        return rows(conn, """SELECT r.id, r.user_id AS userId, r.identifier, r.status, r.created_at AS createdAt,
            r.read_at AS readAt, u.name, u.email, u.department
            FROM password_reset_requests r JOIN users u ON u.id = r.user_id
            WHERE r.status = 'pending' ORDER BY r.created_at DESC""")


@app.post("/api/admin/users/{user_id}/reset-password")
def admin_reset_password(user_id: str, payload: AdminPasswordResetRequest, user: dict[str, Any] = Depends(get_user)) -> dict[str, Any]:
    require_role(user, "admin")
    with closing(connect()) as conn:
        target = row(conn, "SELECT * FROM users WHERE id = ? AND status = 'active'", (user_id,))
        if not target:
            raise HTTPException(404, "员工账号不存在或已停用")
        if target.get("role") == "admin":
            raise HTTPException(403, "不能通过员工重置流程修改管理员密码")
        if payload.request_id:
            request = row(conn, "SELECT id, user_id, status FROM password_reset_requests WHERE id = ?", (payload.request_id,))
            if not request or request.get("user_id") != user_id or request.get("status") != "pending":
                raise HTTPException(409, "密码重置申请无效或已处理")
        stamp = now()
        conn.execute("UPDATE users SET password_hash = ? WHERE id = ?", (hash_password(payload.new_password), user_id))
        conn.execute("DELETE FROM auth_sessions WHERE user_id = ?", (user_id,))
        if payload.request_id:
            conn.execute("UPDATE password_reset_requests SET status = 'handled', handled_by = ?, handled_at = ?, read_at = COALESCE(read_at, ?) WHERE id = ?", (user["id"], stamp, stamp, payload.request_id))
        audit(conn, user["id"], "user.password_reset", "user", user_id, {"request_id": payload.request_id, "target_email": target["email"]})
        conn.commit()
        updated = row(conn, "SELECT id, name, email, department, role, status FROM users WHERE id = ?", (user_id,))
    return updated or {}

@app.get("/api/admin/audit")
def admin_audit(limit: int = 50, user: dict[str, Any] = Depends(get_user)) -> list[dict[str, Any]]:
    require_role(user, "admin")
    safe_limit = max(1, min(limit, 200))
    with closing(connect()) as conn:
        items = rows(conn, "SELECT id, user_id, action, resource_type, resource_id, detail, created_at FROM audit_events ORDER BY created_at DESC LIMIT ?", (safe_limit,))
    for item in items:
        try:
            item["detail"] = json.loads(item["detail"])
        except (TypeError, json.JSONDecodeError):
            item["detail"] = {}
    return items

@app.exception_handler(sqlite3.Error)
def sqlite_error_handler(_, exc: sqlite3.Error) -> JSONResponse:
    return JSONResponse(status_code=500, content={"code": "database_error", "message": str(exc)})
